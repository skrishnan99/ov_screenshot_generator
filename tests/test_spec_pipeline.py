"""Offline end-to-end: base spec -> finished .pptx through the REAL ovdeck.

Only the three model calls are stubbed (assigner, resolver, arranger); the
layout engine, its save() gates, the skeleton store and the canvas invariant
all run for real. Pins the pipeline-order lessons: match before content,
steps numbered after match-skips, required-hole-unmatched skips the slide
with a record, optional absence is silent.

Run: uv run python tests/test_spec_pipeline.py
"""

import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
from _specfix import canned_resolver, keyword_assign, make_run  # noqa: E402

import content as content_mod  # noqa: E402
import matching as matching_mod  # noqa: E402
import arrange as arrange_mod  # noqa: E402
import deckgen  # noqa: E402
import deckspec as ds  # noqa: E402


def main() -> int:
    failures = []
    saved = {
        "assign": matching_mod.assign_call,
        "verify": matching_mod.verify_call,
        "resolve": content_mod.resolve_call,
        "arrange": arrange_mod.arrange_call,
    }
    calls = {"assign": 0, "resolve": 0, "arrange": 0}

    def stub_assign(holes, catalog):
        calls["assign"] += 1
        return keyword_assign(holes, catalog)

    def stub_resolve(reqs, material):
        calls["resolve"] += 1
        return canned_resolver(reqs, material)

    def stub_arrange(title, image_paths, text, feedback="", hint=""):
        calls["arrange"] += 1
        if len(image_paths) == 2:
            # per-side captions must survive into the deck (a real build
            # dropped them: the emitter passed only `caption`)
            return [{"layout": "two_up", "title": title, "images": list(image_paths),
                     "text": {"caption": text.get("text", "")[:80],
                              "left_caption": "LC-CARRIED",
                              "right_caption": "RC-CARRIED"}}]
        slides = []
        for p in image_paths:
            slides.append({"layout": "figure", "title": title, "images": [p],
                           "text": {"caption": text.get("text", "")[:100]}})
        if not slides:
            slides = [{"layout": "statement", "title": title, "images": [],
                       "text": {"intro": "x", "card_title": "S", "bullets": "a\nb\nc"}}]
        return slides

    matching_mod.assign_call = stub_assign
    matching_mod.verify_call = lambda *a, **k: {"match": True, "reason": "stub"}
    content_mod.resolve_call = stub_resolve
    arrange_mod.arrange_call = stub_arrange

    try:
        with tempfile.TemporaryDirectory() as td:
            run = make_run(Path(td))
            out = Path(td) / "out" / "report.pptx"
            out.parent.mkdir(parents=True)

            plan = deckgen.compile_deck(run, out, log=lambda *a: None)

            # ---- the deck exists and passed ovdeck's gates ----
            deck_path = Path(plan.get("deck", ""))
            if not deck_path.exists():
                failures.append("no deck written")
            else:
                from pptx import Presentation

                pr = Presentation(str(deck_path))
                # 2 trained models (Edge Check excluded): title,
                # results_overview, imaging (1), rois (2 holes -> one two_up),
                # training x2 (2 imgs -> one two_up each), results skeleton
                # x2, logic, closing x5  => 14 with this arranger
                if len(pr.slides) < 14:
                    failures.append(f"only {len(pr.slides)} slides in the deck")
                all_text = "\n".join(
                    sh.text_frame.text for sl in pr.slides for sh in sl.shapes
                    if sh.has_text_frame)
                for marker in ("LC-CARRIED", "RC-CARRIED"):
                    if marker not in all_text:
                        failures.append(f"two_up per-side caption {marker} lost")
                kw = pr.core_properties.keywords or ""
                if "ovdeck:template-slides=" not in kw:
                    failures.append("carried-slide stamp missing (audit exemption broken)")

            # ---- one global assignment call (plus at most one repair) ----
            if not (1 <= calls["assign"] <= 2):
                failures.append(f"{calls['assign']} assignment calls; want 1 global (+1 repair)")
            if calls["resolve"] != 1:
                failures.append(f"{calls['resolve']} resolver calls; want exactly 1 batched")

            # ---- steps contiguous in emitted titles ----
            steps = []
            for rec in plan["slides"]:
                t = rec.get("title", "")
                if t.startswith("Step "):
                    steps.append(int(t.split()[1].rstrip(":")))
            if steps != list(range(1, len(steps) + 1)):
                failures.append(f"step numbering not contiguous: {steps}")
            if "{step}" in str(plan["slides"]):
                failures.append("unresolved {step} placeholder reached the plan")

            # ---- matching audit present, optional absence silent ----
            ev = [r for r in plan["slides"] if r["origin"] == "training"]
            if len(ev) != 2:
                failures.append(f"per-model training slides in plan: {len(ev)}, want 2")
            for rec in ev:
                for img in rec["images"]:
                    if img["optional"] and img["path"] is None:
                        pass  # silently absent — exactly right
            if "matching" not in plan or not plan["matching"]:
                failures.append("matching report missing from plan")

            # ---- the new-flow guarantees, end to end ----
            ids = [r["id"] for r in plan["slides"]]
            # ONE results card, after every training slide, before logic
            if ids.count("results") != 1:
                failures.append(f"expected exactly one results card: {ids}")
            else:
                ri = ids.index("results")
                for t in ("training_model-s", "training_horn-quality"):
                    if t not in ids or ids.index(t) > ri:
                        failures.append(f"{t} not before the results card: {ids}")
                if ids.index("logic") < ri:
                    failures.append(f"results card after logic: {ids}")
            if any("edge-check" in i for i in ids):
                failures.append("never-trained model reached the deck")
            # combined-ROI slide matched both trained models' region screens
            rois = next((r for r in plan["slides"] if r["id"] == "rois"), None)
            if rois is None:
                failures.append("rois slide missing from plan")
            else:
                paths = [i["path"] for i in rois["images"]]
                if sorted(p for p in paths if p) != [
                        "deliverables/screenshots/04_roi_horn-quality.png",
                        "deliverables/screenshots/04_roi_model-s.png"]:
                    failures.append(f"rois holes matched {paths}")
            # results_overview got the native raw + overlay pair
            ov = next((r for r in plan["slides"] if r["id"] == "results_overview"), None)
            if ov is None:
                failures.append("results_overview missing from plan")
            else:
                paths = [i["path"] for i in ov["images"]]
                if paths != ["deliverables/images/12_library_raw.jpg",
                             "deliverables/images/12_library_composite.png"]:
                    failures.append(f"results_overview matched {paths}")
            # results card: literal tokens survive to the plan verbatim
            res = next((r for r in plan["slides"] if r["id"] == "results"), None)
            if res is None:
                failures.append("results card missing from plan")
            else:
                toks = res.get("tokens", {})
                if toks.get("train_acc") != "100%" or toks.get("train_imgs") != "6":
                    failures.append(f"results card tokens wrong: {toks}")
            # the library section sits between logic and the closing run
            lib = next((r for r in plan["slides"] if r["id"] == "library"), None)
            if lib is None:
                failures.append("library slide missing from plan")
            else:
                paths = [i["path"] for i in lib["images"]]
                if paths != ["deliverables/screenshots/12_library.png"]:
                    failures.append(f"library hole matched {paths}")
                if not (ids.index("logic") < ids.index("library")
                        < ids.index("closing_capabilities")):
                    failures.append(f"library out of place: {ids}")

            # ---- a required hole that can't match skips the slide, recorded ----
            spec = ds.load_spec()
            spec["slides"].insert(2, {
                "id": "impossible", "layout": "figure", "title": "X",
                "images": [{"expects": "a purple submarine surfacing in a bakery"}],
            })
            spec_path = Path(td) / "mod.yaml"
            import yaml

            spec_path.write_text(yaml.safe_dump(spec, sort_keys=False))
            plan2 = deckgen.compile_deck(run, Path(td) / "out2" / "r.pptx",
                                         spec_path=spec_path, plan_only=True,
                                         log=lambda *a: None)
            skipped = [s for s in plan2["skipped"] if s.get("id") == "impossible"]
            if not skipped or "unmatched" not in skipped[0]["skipped"]:
                failures.append(f"required-unmatched slide not skipped+recorded: {skipped}")
            if any(r["id"] == "impossible" for r in plan2["slides"]):
                failures.append("unmatched slide still compiled")

            # ---- a save-gate rejection triggers ONE re-emit, then succeeds ----
            import ovdeck

            real_save = ovdeck.Deck.save
            state = {"failed": False}

            def flaky_save(self, path=None):
                if not state["failed"]:
                    state["failed"] = True
                    raise ovdeck.LayoutError("synthetic: caption overflow")
                return real_save(self, path)

            ovdeck.Deck.save = flaky_save
            try:
                plan3 = deckgen.compile_deck(run, Path(td) / "out3" / "r.pptx",
                                             log=lambda *a: None)
                if "emit_retried" not in plan3:
                    failures.append("layout-gate retry not recorded in the plan")
                if not Path(plan3.get("deck", "")).exists():
                    failures.append("retry did not produce a deck")
                ids3 = [r["id"] for r in plan3["slides"]]
                if len(ids3) != len(set(ids3)):
                    failures.append("retry duplicated slide records in the plan")
            finally:
                ovdeck.Deck.save = real_save
    finally:
        matching_mod.assign_call = saved["assign"]
        matching_mod.verify_call = saved["verify"]
        content_mod.resolve_call = saved["resolve"]
        arrange_mod.arrange_call = saved["arrange"]

    if failures:
        for f in failures:
            print(f"FAIL: {f}")
        return 1
    print("ALL SPEC-PIPELINE CHECKS PASSED")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
