"""Standing slide copy survives arrangement verbatim, and renders as the
split's own purple intro section.

The imaging slide carries a literal `intro` token — the spec author's exact
words, every deck. Carrying text used to be prompt-enforced only; literal
tokens are now MUST-CARRY, validated in code. What this suite pins:

- the default spec's imaging slide carries the literal intro (purple
  framing copy) ahead of the dynamic text,
- validate_arrangement rejects an arrangement that dropped or paraphrased
  a literal, and accepts verbatim carriage on any slide/field,
- the arrange ladder retries with the missing literal named, and the
  deterministic fallback preserves literals by construction (split intro
  for one image — never a truncated figure caption — and an untruncated
  caption lead otherwise),
- ovdeck's split renders the intro in ACCENT purple between the card
  title and the body content,
- deckgen passes literal tokens (and only literals) as must-carry.

Run: uv run python tests/test_static_intro.py
"""

import inspect
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import _specfix  # noqa: F401, E402  (sys.path + data-dir isolation)

import yaml  # noqa: E402

import arrange as arrange_mod  # noqa: E402

REPO = Path(__file__).resolve().parent.parent
SPEC = REPO / "skills" / "overview-deck" / "specs" / "default-deck.yaml"

LIT = ("In the image setup page you define the best configuration for "
       "your application — exposure, gamma and more. It's one of the most "
       "important steps, as everything else will be defined by this.")


def main() -> int:
    failures = []

    # ---- the spec carries the literal, on the imaging slide ----
    spec = yaml.safe_load(SPEC.read_text())
    imaging = next(s for s in spec["slides"] if s.get("id") == "imaging")
    intro = imaging.get("tokens", {}).get("intro")
    if not isinstance(intro, str) or "exposure, gamma and more" not in intro:
        failures.append(f"imaging intro literal missing/changed: {intro!r}")
    if "intro" not in imaging.get("hint", "").lower():
        failures.append("imaging hint lost the intro placement instruction")
    # the section heading between intro and dynamic text: a literal, so it
    # rides the same must-carry guarantee
    if imaging.get("tokens", {}).get("heading") != "Capture Configuration":
        failures.append(f"imaging heading literal: "
                        f"{imaging.get('tokens', {}).get('heading')!r}")
    if "[heading]" not in imaging.get("hint", ""):
        failures.append("imaging hint lost the heading mapping")
    # the aligner slide's heading is dynamic (like the logic slide's),
    # naming the aligner in use; the slide itself still drops when the
    # aligner is skipped
    aligner = next(s for s in spec["slides"] if s.get("id") == "aligner")
    ah = aligner.get("tokens", {}).get("heading")
    if not (isinstance(ah, dict) and "llm" in ah):
        failures.append(f"aligner heading is not a resolver hole: {ah!r}")
    else:
        for phrase in ("Deep Learning Alignment", "Classical Alignment"):
            if phrase not in ah["llm"]:
                failures.append(f"aligner heading guidance lost {phrase!r}")
    if aligner.get("when") != {"aligner.skipped": False}:
        failures.append(f"aligner skip condition changed: {aligner.get('when')}")

    # ---- validation: dropped or paraphrased literal is rejected ----
    ok_plan = [{"layout": "split", "title": "T", "images": ["a.png"],
                "text": {"card_title": "S", "intro": LIT, "para": "dynamic"}}]
    if arrange_mod.validate_arrangement(ok_plan, ["a.png"], (LIT,)):
        failures.append("verbatim carriage was rejected")
    dropped = [{"layout": "figure", "title": "T", "images": ["a.png"],
                "text": {"caption": "dynamic only"}}]
    probs = arrange_mod.validate_arrangement(dropped, ["a.png"], (LIT,))
    if not any("literal text dropped" in p for p in probs):
        failures.append(f"dropped literal not flagged: {probs}")
    para = [{"layout": "split", "title": "T", "images": ["a.png"],
             "text": {"card_title": "S",
                      "intro": "Imaging setup is important.", "para": "x"}}]
    if not arrange_mod.validate_arrangement(para, ["a.png"], (LIT,)):
        failures.append("paraphrased literal accepted as carried")
    # whitespace differences are not violations
    rewrapped = [{"layout": "split", "title": "T", "images": ["a.png"],
                  "text": {"card_title": "S",
                           "intro": " ".join(LIT.split()[:9]) + "\n"
                                    + " ".join(LIT.split()[9:]),
                           "para": "x"}}]
    if arrange_mod.validate_arrangement(rewrapped, ["a.png"], (LIT,)):
        failures.append("rewrapped (whitespace-only) literal rejected")

    # ---- the ladder: drop -> retry named -> success; twice -> fallback ----
    calls = []

    def flaky(title, image_paths, text, feedback="", hint=""):
        calls.append(feedback)
        if len(calls) == 1:
            return dropped
        return ok_plan

    saved = arrange_mod.arrange_call
    try:
        arrange_mod.arrange_call = flaky
        plan = arrange_mod.arrange("T", ["a.png"], {"intro": LIT, "text": "dynamic"},
                                   must_carry=(LIT,), log=lambda *a: None)
        if plan != ok_plan:
            failures.append("retry did not recover the literal")
        if len(calls) != 2 or "literal text dropped" not in calls[1]:
            failures.append(f"retry feedback did not name the literal: {calls}")

        calls.clear()
        arrange_mod.arrange_call = lambda *a, **k: dropped
        plan = arrange_mod.arrange("T", ["a.png"], {"intro": LIT, "text": "dynamic"},
                                   must_carry=(LIT,), log=lambda *a: None)
        got = " ".join(v for s in plan for v in s["text"].values())
        if " ".join(LIT.split()) not in " ".join(got.split()):
            failures.append("fallback lost the literal")
        if plan[0]["layout"] != "split" or plan[0]["text"].get("intro") != LIT:
            failures.append(f"one-image fallback should be a split intro: {plan}")
    finally:
        arrange_mod.arrange_call = saved

    # ---- fallback with two images: literal leads the caption untruncated ----
    plan = arrange_mod.fallback_arrangement("T", ["a.png", "b.png"],
                                            {"intro": LIT, "text": "d" * 300},
                                            must_carry=(LIT,))
    cap = plan[0]["text"].get("caption", "")
    if not cap.startswith(LIT):
        failures.append(f"two-image fallback truncated the literal: {cap[:80]!r}")

    # ---- ovdeck: intro renders purple, between card title and body ----
    from ovdeck import ACCENT, Deck

    with tempfile.TemporaryDirectory() as td:
        from PIL import Image

        img = Path(td) / "s.png"
        Image.new("RGB", (640, 400), (40, 40, 60)).save(img)
        out = Path(td) / "d.pptx"
        # strict=False only to skip the deck-level opening/closing logo
        # rule, which a one-slide unit deck cannot satisfy; the slide's own
        # layout checks are asserted explicitly below.
        d = Deck(str(out), strict=False)
        d.split("T", str(img), card_title="Card", intro=LIT,
                bullets=["one: 1", "two: 2"])
        slide_issues = [i for i in d.check() if "missing-logo" not in str(i)]
        if slide_issues:
            failures.append(f"intro split fails layout checks: {slide_issues}")
        d.save()

        from pptx import Presentation

        prs = Presentation(str(out))
        boxes = {}
        for sh in prs.slides[-1].shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                boxes[sh.text_frame.text.strip()[:20]] = sh
        intro_sh = next((sh for k, sh in boxes.items() if k.startswith("In the image")), None)
        card_sh = next((sh for k, sh in boxes.items() if k == "Card"), None)
        bullet_sh = next((sh for k, sh in boxes.items() if "one: 1" in sh.text_frame.text), None)
        if not (intro_sh and card_sh and bullet_sh):
            failures.append(f"split shapes missing: {list(boxes)}")
        else:
            # the intro TOPS the card; the card title renders below it as a
            # section heading over the body (the mockup's order)
            if not (intro_sh.top < card_sh.top < bullet_sh.top):
                failures.append("card order is not intro -> heading -> body")
            run = intro_sh.text_frame.paragraphs[0].runs[0]
            if run.font.color.rgb != ACCENT:
                failures.append(f"intro colour {run.font.color.rgb} != accent {ACCENT}")

    # ---- deckgen passes literals (and only literals) as must-carry ----
    src = inspect.getsource(__import__("deckgen"))
    if "must_carry=carry" not in src:
        failures.append("deckgen no longer passes literal tokens as must-carry")

    if failures:
        print("FAILURES:")
        for f in failures:
            print(" -", f)
        return 1
    print("ALL STATIC-INTRO CHECKS PASSED")
    return 0


if __name__ == "__main__":
    sys.exit(main())
