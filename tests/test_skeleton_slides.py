"""Boilerplate slides come from owned skeletons, filled and placed exactly.

The library/capabilities/team/thank-you slides used to be re-authored by the
layout engine (a different approximation every build) and then transplanted
verbatim from the 26 MB reference template (byte-exact but unownable — its
defects shipped into every deck). They are now extracted ONCE into small
single-slide skeletons the skill owns, with sidecar YAML describing content
holes; builds fill the holes and append the slide unchanged.

What this pins:
  - the skeleton store is complete, small, and fixed: no template page
    numbers, the library subtitle has wrap headroom (the defect that started
    all this), media re-embedded at display size
  - holes are placeholders, not pictures: filling targets the "Insert
    screenshot here" shape and never replaces standing content
  - an unfilled hole raises rather than shipping placeholder text
  - the ovdeck integration stamps carried-slide indices so the audit knows
    what the engine did not lay out

Run: uv run python tests/test_skeleton_slides.py
"""

import io
import re
import sys
from pathlib import Path

REPO = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(REPO))
SCRIPTS = REPO / "skills" / "overview-deck" / "scripts"
sys.path.insert(0, str(SCRIPTS))

import template_slides as ts  # noqa: E402

PAGE_RE = re.compile(r"^\s*\d{1,2}\s*/\s*\d{1,2}\s*$")


def _png(tmp: Path, size=(1600, 1000)) -> Path:
    from PIL import Image

    p = tmp / "shot.png"
    Image.new("RGB", size, (40, 40, 60)).save(p)
    return p


def main() -> int:
    import tempfile

    from pptx import Presentation
    from pptx.util import Emu

    failures = []

    # ---- store: complete, small, single-slide, fixed ----
    total = 0
    for name in ts.TEMPLATE_SLIDES:
        try:
            p = ts.skeleton_path(name)
        except ts.TemplateError as e:
            failures.append(f"missing skeleton: {e}")
            continue
        size = p.stat().st_size
        total += size
        if size > 1_500_000:
            failures.append(f"{name}.pptx is {size/1024:.0f} KB; media shrink regressed")
        pres = Presentation(str(p))
        if len(pres.slides) != 1:
            failures.append(f"{name}.pptx has {len(pres.slides)} slides, want 1")
            continue
        slide = pres.slides[0]
        pages = [sh.text_frame.text for sh in slide.shapes
                 if sh.has_text_frame and PAGE_RE.match(sh.text_frame.text or "")]
        if pages:
            failures.append(f"{name}: template page number survived: {pages}")
    if total > 3_500_000:
        failures.append(f"skeleton store is {total/1048576:.1f} MB; want < 3.5 MB")

    # THE CANVAS INVARIANT: every skeleton must be on ovdeck's canvas, or it
    # renders huddled in the top-left of the deck (absolute EMUs on a bigger
    # slide). This is the defect that prompted the normalisation.
    from pptx.util import Inches

    from ovdeck import SLIDE_H, SLIDE_W

    for name in ts.TEMPLATE_SLIDES:
        pres = Presentation(str(ts.skeleton_path(name)))
        if (pres.slide_width, pres.slide_height) != (Inches(SLIDE_W), Inches(SLIDE_H)):
            failures.append(
                f"{name}: canvas {pres.slide_width}x{pres.slide_height} EMU != "
                f"deck canvas {SLIDE_W}x{SLIDE_H} in"
            )

    # Anti-shadowing: _shrink_pictures once reused its `scale` parameter as
    # the per-picture resize ratio, compounding down the slide (rasters of
    # 687 -> 180 -> 47 -> 12 -> 3 px). Healthy skeletons have no tiny rasters
    # and large frames keep detailed rasters.
    import io as _io

    from pptx.oxml.ns import qn
    from PIL import Image as _Image

    cap = Presentation(str(ts.skeleton_path("capabilities"))).slides[0]
    for sh in cap.shapes:
        if sh.shape_type != 13:
            continue
        part = sh.part.related_part(
            sh._element.blipFill.find(qn("a:blip")).get(qn("r:embed"))
        )
        with _Image.open(_io.BytesIO(part.blob)) as im:
            w_in = Emu(sh.width).inches
            if min(im.size) < 90:
                failures.append(
                    f"capabilities: {w_in:.2f}in picture raster {im.size} — "
                    f"compounding shrink is back"
                )
            if w_in > 3 and im.width < 500:
                failures.append(
                    f"capabilities: {w_in:.2f}in frame with {im.width}px raster"
                )

    # the fix that started all this: subtitle wrap headroom on the library
    lib = Presentation(str(ts.skeleton_path("library"))).slides[0]
    sub = [Emu(sh.height).inches for sh in lib.shapes
           if sh.has_text_frame and sh.text_frame.text.startswith("Easier root cause")]
    if not sub or sub[0] < 0.55:
        failures.append(f"library subtitle headroom fix missing (H={sub})")

    # ---- the type scale stays inside the skeleton-grounded ranges ----
    # Authored slides must read at the same size as the standing template
    # slides: body there is 13-18.7pt, card titles 15-18, slide titles 28-48,
    # floor 11 (the dense team/integration slides). A scale that drifts back
    # down makes generated content illegibly small next to the carried slides.
    import json as _json

    tokens = _json.loads(
        (REPO / "skills" / "overview-deck" / "assets" / "tokens.json").read_text()
    )["typography"]
    sz, rng = tokens["scale_pt"], tokens.get("range_pt") or {}
    if not rng:
        failures.append("tokens.typography.range_pt missing — the grounded ranges are gone")
    floor = rng.get("floor", 11)
    for k, v in sz.items():
        if v < floor:
            failures.append(f"scale_pt[{k}] = {v} is below the {floor}pt floor")
    lo, hi = rng.get("body", [13, 19])
    if not (lo <= sz["body"] <= hi and lo <= sz["bullet"] + 1):
        failures.append(f"body/bullet ({sz['body']}/{sz['bullet']}) outside body range {lo}-{hi}")
    lo, hi = rng.get("card_titles", [15, 18])
    if not (lo <= sz["card_title"] <= hi):
        failures.append(f"card_title {sz['card_title']} outside {lo}-{hi}")
    lo, hi = rng.get("slide_titles", [28, 48])
    if not (lo <= sz["slide_title"] <= hi):
        failures.append(f"slide_title {sz['slide_title']} outside {lo}-{hi}")

    # ---- query: profile sees the hole, sidecar names it, no drift ----
    prof = ts.profile("library")
    if [s["name"] for s in prof["slots"]] != ["library_screen"]:
        failures.append(f"library profile slots: {[s['name'] for s in prof['slots']]}")
    if prof["warnings"]:
        failures.append(f"library sidecar drift: {prof['warnings']}")
    if not prof["slots"] or prof["slots"][0]["is_picture"]:
        failures.append("library hole should be a placeholder, not a picture")

    # ---- fill + place ----
    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)
        shot = _png(tmp)

        def base():
            src = Presentation(str(ts.skeleton_path("thank_you")))
            b = Presentation()
            b.slide_width, b.slide_height = src.slide_width, src.slide_height
            return b

        # library filled: picture lands inside the placeholder box, note gone
        b = base()
        ts.append(b, "library", image=shot)
        slide = Presentation(io.BytesIO(_save(b))).slides[0]
        pics = [sh for sh in slide.shapes if sh.shape_type == 13]
        if len(pics) != 1:
            failures.append(f"library fill produced {len(pics)} pictures, want 1")
        else:
            # The hole's frame, read from the skeleton itself so the test
            # tracks the canvas normalisation instead of hardcoding EMUs.
            from deck.assemble import image_slots as _slots

            hole = _slots(Presentation(str(ts.skeleton_path("library"))).slides[0])[0]
            box = (Emu(hole.left).inches, Emu(hole.top).inches,
                   Emu(hole.width).inches, Emu(hole.height).inches)
            L, T = Emu(pics[0].left).inches, Emu(pics[0].top).inches
            W, H = Emu(pics[0].width).inches, Emu(pics[0].height).inches
            if not (box[0] - 0.02 <= L and T >= box[1] - 0.02
                    and L + W <= box[0] + box[2] + 0.02
                    and T + H <= box[1] + box[3] + 0.02):
                failures.append(f"library image outside placeholder box: {(L, T, W, H)}")
        texts = " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
        if "insert screenshot" in texts.lower():
            failures.append("placeholder note shipped on the filled library slide")

        # unfilled hole must raise, not ship placeholder text
        try:
            ts.append(base(), "library")
            failures.append("library without an image should raise")
        except ts.TemplateError:
            pass

        # a skeleton with no hole must reject an image, and keep its pictures
        try:
            ts.append(base(), "capabilities", image=shot)
            failures.append("capabilities has no hole; image should raise")
        except ts.TemplateError:
            pass
        b = base()
        ts.append(b, "capabilities")
        slide = Presentation(io.BytesIO(_save(b))).slides[0]
        n_pics = sum(1 for sh in slide.shapes if sh.shape_type == 13)
        if n_pics < 5:
            failures.append(f"capabilities standing pictures lost ({n_pics} left)")

        # unknown names and surplus tokens are errors
        for kwargs, label in ((dict(), "unknown name"),
                              (dict(tokens={"nope": "x"}), "surplus token")):
            try:
                if "tokens" in kwargs:
                    ts.append(base(), "team", **kwargs)
                else:
                    ts.append(base(), "not_a_slide")
                failures.append(f"{label} should raise")
            except ts.TemplateError:
                pass

        # ---- ovdeck integration: stamp + carried count + alias ----
        from ovdeck import Deck

        d = Deck(str(tmp / "deck.pptx"))
        d.title_slide("OV80i", "Test")
        d.skeleton_slide("library", image=shot)
        d.template_slide("thank_you")  # deprecated alias still works
        out = d.save()
        pr = Presentation(str(out))
        if len(pr.slides) != 3:
            failures.append(f"deck has {len(pr.slides)} slides, want 3")
        kw = pr.core_properties.keywords or ""
        if "ovdeck:template-slides=2,3" not in kw:
            failures.append(f"carried-slide stamp wrong: {kw!r}")

    if failures:
        for f in failures:
            print(f"FAIL: {f}")
        return 1
    print("ALL SKELETON-SLIDE CHECKS PASSED")
    return 0


def _save(prs) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


if __name__ == "__main__":
    raise SystemExit(main())
