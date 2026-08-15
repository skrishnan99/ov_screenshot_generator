"""The contact slide signs with the SE who ran the visit, not the template's
static person.

thank_you.pptx's contact block is tokenized at extraction ({{contact_*}})
and filled at the lowest level — template_slides.append — from the per-user
profile at data_dir()/engineer.json, so every caller (the v2 closing run,
hand-written v1 scripts) gets the right contact without passing anything.
What this suite pins:

- the skeleton carries exactly the three contact tokens,
- the profile loader: full profile, field-wise partial, env override wins,
  malformed file and missing file both degrade to the visibly generic
  placeholders ("SE Name", ...) — never an exception,
- append() fills the slide from the profile, and from placeholders when
  none exists; no {{token}} ever survives to the deck,
- compile_deck records the contact source in the plan so a placeholder
  contact is surfaced, never shipped unnoticed.

Run: uv run python tests/test_engineer_contact.py
"""

import json
import os
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import _specfix  # noqa: E402  (isolates OV_REPORT_DATA_DIR first)
from _specfix import make_run  # noqa: E402

REPO = Path(__file__).resolve().parent.parent


def _slide_text(prs) -> str:
    out = []
    for sh in prs.slides[0].shapes:
        if sh.has_text_frame:
            out.append(sh.text_frame.text)
    return "\n".join(out)


def _fresh_data_dir():
    d = tempfile.mkdtemp(prefix="sg-eng-")
    os.environ["OV_REPORT_DATA_DIR"] = d
    return Path(d)


def main() -> int:
    failures = []

    from pptx import Presentation
    from pptx.util import Inches

    import template_slides as ts
    from core import engineer
    from deck.assemble import find_tokens

    # ---- the skeleton carries exactly the three contact tokens ----
    skel = Presentation(str(ts.skeleton_path("thank_you")))
    toks = set(find_tokens(skel.slides[0]))
    if toks != {"contact_name", "contact_email", "contact_phone"}:
        failures.append(f"thank_you tokens: {sorted(toks)}")

    # ---- loader: missing file -> placeholders ----
    _fresh_data_dir()
    contact, source = engineer.load_profile()
    if source != "placeholder" or contact["name"] != "SE Name" \
            or contact["email"] != "SE Email" \
            or contact["phone"] != "SE Contact Number":
        failures.append(f"missing-profile load: {source} {contact}")

    # ---- loader: full profile ----
    d = _fresh_data_dir()
    engineer.save_profile("Jane Doe", "jane@overview.ai", "(555) 000-1111")
    contact, source = engineer.load_profile()
    if source != "profile" or contact["name"] != "Jane Doe":
        failures.append(f"full-profile load: {source} {contact}")

    # ---- phone renders in US display format, from any stored shape ----
    for raw, want in (
        ("9096156153", "(909) 615-6153"),          # bare digits
        ("(909) 615-6153", "(909) 615-6153"),      # already formatted: stable
        ("909-615-6153", "(909) 615-6153"),        # dashed
        ("1 909 615 6153", "(909) 615-6153"),      # 11 with leading 1
        ("+44 20 7946 0958", "+44 20 7946 0958"),  # non-US: untouched
        ("x1234", "x1234"),                        # short/odd: untouched
    ):
        got = engineer.format_phone(raw)
        if got != want:
            failures.append(f"format_phone({raw!r}) -> {got!r}, want {want!r}")
    engineer.save_profile("Jane Doe", "jane@overview.ai", "9096156153")
    contact, _ = engineer.load_profile()
    if contact["phone"] != "(909) 615-6153":
        failures.append(f"loaded phone not display-formatted: {contact['phone']!r}")
    if json.loads((d / "engineer.json").read_text())["phone"] != "9096156153":
        failures.append("stored phone was rewritten; formatting must be display-time only")

    # ---- loader: partial file is field-wise ----
    (d / "engineer.json").write_text(json.dumps({"name": "Jane Doe"}))
    contact, source = engineer.load_profile()
    if source != "partial" or contact["name"] != "Jane Doe" \
            or contact["email"] != "SE Email":
        failures.append(f"partial-profile load: {source} {contact}")

    # ---- loader: env override wins over the file ----
    os.environ["SG_ENGINEER_NAME"] = "Env Person"
    try:
        contact, _ = engineer.load_profile()
        if contact["name"] != "Env Person":
            failures.append(f"env override lost: {contact['name']!r}")
    finally:
        del os.environ["SG_ENGINEER_NAME"]

    # ---- loader: malformed file degrades, never raises ----
    (d / "engineer.json").write_text("{not json")
    try:
        contact, source = engineer.load_profile()
        if source != "placeholder":
            failures.append(f"malformed profile source: {source}")
    except Exception as e:
        failures.append(f"malformed profile raised: {e}")

    # ---- missing_fields: names exactly what the mandatory ask must
    # request; env overrides count as present ----
    d = _fresh_data_dir()
    if engineer.missing_fields() != ["name", "email", "phone"]:
        failures.append(f"empty profile missing_fields: {engineer.missing_fields()}")
    (d / "engineer.json").write_text(json.dumps(
        {"name": "Jane Doe", "email": "jane@overview.ai"}))
    if engineer.missing_fields() != ["phone"]:
        failures.append(f"partial profile missing_fields: {engineer.missing_fields()}")
    os.environ["SG_ENGINEER_PHONE"] = "9096156153"
    try:
        if engineer.missing_fields():
            failures.append("env-satisfied field still reported missing")
    finally:
        del os.environ["SG_ENGINEER_PHONE"]

    # ---- save_profile merges: filling ONE missing field keeps the rest ----
    engineer.save_profile(phone="9096156153")
    stored = json.loads((d / "engineer.json").read_text())
    if stored != {"name": "Jane Doe", "email": "jane@overview.ai",
                  "phone": "9096156153"}:
        failures.append(f"single-field save clobbered the profile: {stored}")
    if engineer.missing_fields():
        failures.append(f"profile still incomplete after merge save")

    # ---- preflight: hard check under --ensure-engineer-profile only ----
    import subprocess
    env = dict(os.environ)
    env["OV_REPORT_DATA_DIR"] = _fresh_data_dir().as_posix()
    base = [sys.executable, str(REPO / "preflight.py")]
    r = subprocess.run(base + ["--ensure-engineer-profile"],
                       capture_output=True, text=True, env=env, cwd=REPO)
    if r.returncode == 0:
        failures.append("preflight passed with an empty profile under the flag")
    if "missing: name, email, phone" not in r.stdout:
        failures.append(f"preflight does not name the missing fields:\n{r.stdout[-300:]}")
    r = subprocess.run(base, capture_output=True, text=True, env=env, cwd=REPO)
    if r.returncode != 0:
        failures.append(f"preflight failed WITHOUT the flag:\n{r.stdout[-300:]}")
    (Path(env["OV_REPORT_DATA_DIR"]) / "engineer.json").write_text(json.dumps(
        {"name": "Jane Doe", "email": "jane@overview.ai", "phone": "9096156153"}))
    r = subprocess.run(base + ["--ensure-engineer-profile"],
                       capture_output=True, text=True, env=env, cwd=REPO)
    if r.returncode != 0:
        failures.append(f"preflight failed with a complete profile:\n{r.stdout[-300:]}")

    # ---- append(): placeholder fill when no profile exists ----
    from ovdeck import SLIDE_H, SLIDE_W

    def _fill():
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)
        ts.append(prs, "thank_you")
        # a transplanted slide is a raw OPC part in memory — save and
        # reopen before the typed shape API can read it
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs.save(f.name)
            return _slide_text(Presentation(f.name))

    _fresh_data_dir()
    text = _fill()
    for want in ("SE Name", "SE Email", "SE Contact Number"):
        if want not in text:
            failures.append(f"placeholder fill missing {want!r}")
    if "{{" in text:
        failures.append("unfilled token survived placeholder fill")

    # ---- append(): profile fill ----
    _fresh_data_dir()
    engineer.save_profile("Jane Doe", "jane@overview.ai", "(555) 000-1111")
    text = _fill()
    if "Jane Doe" not in text or "jane@overview.ai" not in text:
        failures.append(f"profile fill wrong: {text[:120]!r}")
    if "SE Name" in text or "{{" in text:
        failures.append("profile fill left placeholders or tokens")

    # ---- compile_deck records the contact source in the plan ----
    import deckspec as ds  # noqa: F401  (sys.path via _specfix)
    import matching as matching_mod
    from deckgen import compile_deck

    saved = (matching_mod.assign_call, matching_mod.verify_call,
             matching_mod.block_quality_call)
    try:
        matching_mod.assign_call = lambda holes, catalog: [
            {"hole": h.id, "path": None, "confidence": "high", "reason": "s"}
            for h in holes]
        matching_mod.verify_call = lambda *a, **k: {"match": True, "reason": "s"}
        matching_mod.block_quality_call = lambda desc: {
            "product_image": True, "annotated": True, "reason": "s"}
        with tempfile.TemporaryDirectory() as td:
            run = make_run(Path(td))
            _fresh_data_dir()
            plan = compile_deck(run, Path(td) / "out.pptx", plan_only=True,
                                log=lambda *a: None)
            if plan.get("contact", {}).get("source") != "placeholder":
                failures.append(f"plan contact (no profile): {plan.get('contact')}")
            engineer.save_profile("Jane Doe", "jane@overview.ai", "(555) 000-1111")
            plan = compile_deck(run, Path(td) / "out.pptx", plan_only=True,
                                log=lambda *a: None)
            if plan.get("contact") != {"source": "profile", "name": "Jane Doe"}:
                failures.append(f"plan contact (profile): {plan.get('contact')}")
    finally:
        (matching_mod.assign_call, matching_mod.verify_call,
         matching_mod.block_quality_call) = saved

    if failures:
        print("FAILURES:")
        for f in failures:
            print(" -", f)
        return 1
    print("ALL ENGINEER-CONTACT CHECKS PASSED")
    return 0


if __name__ == "__main__":
    sys.exit(main())
