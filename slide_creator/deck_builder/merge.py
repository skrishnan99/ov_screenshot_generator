"""Merge single-slide .pptx files into one deck.

Copied from ``recipe_decryption/case_study/preview/merge.py`` (minus
the PDF-splitting helper this pipeline doesn't need). The non-obvious
parts — effective-background copying, theme-reference baking, and
placeholder text-color baking — exist because each skeleton was
exported from a different Google Slides source deck; without them the
merged deck re-resolves scheme colors/fonts against the wrong theme
and dark designs render with invisible text.

Dependencies: python-pptx, lxml.
"""

from __future__ import annotations

from pathlib import Path
from typing import Sequence


def merge_pptx(pptx_paths: Sequence[Path], output_path: Path) -> Path:
    """Merge multiple single-slide .pptx files into one multi-slide .pptx.

    Each source file contributes its first slide. Slide dimensions are
    taken from the first file. Shapes, images, and text are copied
    via the low-level lxml clone so formatting, fonts, and embedded
    images survive the merge.

    Args:
        pptx_paths: Ordered list of .pptx files, each containing one slide.
        output_path: Where to write the merged .pptx.

    Returns:
        ``output_path``.
    """
    from pptx import Presentation
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from lxml import etree
    from copy import deepcopy
    import io

    if not pptx_paths:
        raise ValueError("merge_pptx: empty pptx_paths list")

    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Start with the first file as the base — this preserves its
    # slide dimensions, theme, and master layouts.
    merged = Presentation(str(pptx_paths[0]))

    for slide_idx, src_path in enumerate(pptx_paths[1:], start=1):
        src_prs = Presentation(str(src_path))
        if not src_prs.slides:
            continue
        src_slide = src_prs.slides[0]

        # Add a blank slide using the merged presentation's first layout.
        layout = merged.slide_layouts[6]  # blank layout
        new_slide = merged.slides.add_slide(layout)

        # Copy the EFFECTIVE background from the source: the slide's
        # own <p:bg>, else the layout's, else the master's — Google
        # Slides designs usually put the background on the layout or
        # master, which the merged deck's blank layout doesn't have.
        # Without this, dark designs render white and their light text
        # becomes invisible.
        ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        new_cSld = new_slide._element.find(f"{{{ns_p}}}cSld")
        src_bg, bg_owner_part = _effective_background(src_slide)
        if src_bg is not None and new_cSld is not None:
            old_bg = new_cSld.find(f"{{{ns_p}}}bg")
            if old_bg is not None:
                new_cSld.remove(old_bg)
            bg_clone = deepcopy(src_bg)
            # Background image fills reference the owning part's rels —
            # re-home them onto the new slide.
            _rehome_blips(bg_clone, bg_owner_part, new_slide, slide_idx, merged)
            new_cSld.insert(0, bg_clone)

        # Clear the blank slide's shape tree.
        sp_tree = new_slide.shapes._spTree
        for child in list(sp_tree):
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag in ("sp", "pic", "graphicFrame", "grpSp", "cxnSp"):
                sp_tree.remove(child)

        # Build a mapping from source rIds to new rIds so we can
        # rewrite <a:blip r:embed="rIdN"> references in cloned shapes.
        rid_map: dict[str, str] = {}

        for rel in src_slide.part.rels.values():
            if "image" in rel.reltype:
                # Read the image bytes from the source part and create
                # a fresh image part in the merged slide. This avoids
                # the "Duplicate name: ppt/media/imageN.png" collision
                # because python-pptx assigns a unique partname.
                image_bytes = rel.target_part.blob
                content_type = rel.target_part.content_type

                from pptx.opc.package import Part
                from pptx.opc.packuri import PackURI

                # Generate a unique partname using the slide index.
                ext = {
                    "image/png": ".png",
                    "image/jpeg": ".jpeg",
                    "image/jpg": ".jpg",
                    "image/gif": ".gif",
                    "image/svg+xml": ".svg",
                }.get(content_type, ".png")
                unique_name = f"/ppt/media/slide{slide_idx}_img{len(rid_map)}{ext}"

                image_part = Part(
                    PackURI(unique_name),
                    content_type,
                    blob=image_bytes,
                    package=merged.part.package,
                )
                new_rid = new_slide.part.relate_to(image_part, rel.reltype)
                rid_map[rel.rId] = new_rid

        # Clone each shape, rewriting image rId references.
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
              "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships"}

        src_sp_tree = src_slide.shapes._spTree
        for child in src_sp_tree:
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag in ("sp", "pic", "graphicFrame", "grpSp", "cxnSp"):
                cloned = deepcopy(child)

                # Rewrite r:embed attributes on <a:blip> elements.
                for blip in cloned.iter("{%s}blip" % ns["a"]):
                    old_rid = blip.get("{%s}embed" % ns["r"])
                    if old_rid and old_rid in rid_map:
                        blip.set("{%s}embed" % ns["r"], rid_map[old_rid])

                sp_tree.append(cloned)

        # Bake the SOURCE deck's theme into the copied slide. The
        # merged deck keeps only the base file's theme, so any scheme
        # color (accent1, bg1, ...) or theme font (+mn-lt, ...) in the
        # copied XML would re-resolve against the WRONG theme — every
        # skeleton exported from a different Google Slides source deck
        # drifted visibly. Resolving those references to literal values
        # makes each slide theme-independent.
        color_map, font_map = _slide_theme_maps(src_slide)
        _bake_theme_references(new_slide._element, color_map, font_map)

        # Placeholder text (titles etc.) often carries NO explicit
        # color — it inherits from the source deck's layout/master
        # text styles, which don't travel with the slide. Resolve the
        # inherited color and write it onto colorless runs so the
        # merged deck can't restyle them.
        _bake_placeholder_text_colors(new_slide._element, src_slide, color_map)

    merged.save(str(output_path))
    return output_path


# ---------------------------------------------------------------------------
# Theme baking — make copied slides independent of the base theme
# ---------------------------------------------------------------------------

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _effective_background(src_slide):
    """The background that actually renders: slide → layout → master.

    Returns ``(bg_element, owning_part)`` or ``(None, None)``.
    """
    candidates = (
        (src_slide._element, src_slide.part),
        (src_slide.slide_layout._element, src_slide.slide_layout.part),
        (src_slide.slide_layout.slide_master._element,
         src_slide.slide_layout.slide_master.part),
    )
    for element, part in candidates:
        cSld = element.find(f"{{{_P_NS}}}cSld")
        if cSld is None:
            continue
        bg = cSld.find(f"{{{_P_NS}}}bg")
        if bg is not None:
            return bg, part
    return None, None


def _rehome_blips(element, owner_part, new_slide, slide_idx, merged) -> None:
    """Copy image parts referenced by ``element`` onto the new slide.

    Cloned XML (e.g. a background picture fill) carries r:embed ids
    that only resolve against the SOURCE part's relationships; without
    re-homing they dangle and the package is corrupt.
    """
    if owner_part is None:
        return
    from pptx.opc.package import Part
    from pptx.opc.packuri import PackURI

    for i, blip in enumerate(element.iter(f"{{{_A_NS}}}blip")):
        rid = blip.get(f"{{{_R_NS}}}embed")
        if not rid:
            continue
        try:
            target = owner_part.rels[rid].target_part
        except KeyError:
            continue
        ext = {
            "image/png": ".png", "image/jpeg": ".jpeg",
            "image/jpg": ".jpg", "image/gif": ".gif",
        }.get(target.content_type, ".png")
        part = Part(
            PackURI(f"/ppt/media/slide{slide_idx}_bg{i}{ext}"),
            target.content_type, blob=target.blob,
            package=merged.part.package,
        )
        new_rid = new_slide.part.relate_to(
            part, owner_part.rels[rid].reltype
        )
        blip.set(f"{{{_R_NS}}}embed", new_rid)


def _slide_theme_maps(src_slide) -> tuple[dict, dict]:
    """Resolve a slide's theme into literal lookup tables.

    Returns ``(colors, fonts)``: scheme-color name → hex (including
    the master's bg1/tx1 aliases) and theme-font token ("+mn-lt", …)
    → typeface name. Empty dicts when anything is missing — baking
    then no-ops and the slide renders as before.
    """
    from lxml import etree
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    try:
        master = src_slide.slide_layout.slide_master
        theme_root = etree.fromstring(master.part.part_related_by(RT.THEME).blob)

        colors: dict = {}
        clr_scheme = theme_root.find(f".//{{{_A_NS}}}clrScheme")
        for child in clr_scheme if clr_scheme is not None else []:
            name = etree.QName(child).localname
            srgb = child.find(f"{{{_A_NS}}}srgbClr")
            sysc = child.find(f"{{{_A_NS}}}sysClr")
            val = (
                srgb.get("val") if srgb is not None
                else sysc.get("lastClr") if sysc is not None
                else None
            )
            if val:
                colors[name] = val

        # Color aliases (bg1/tx1/bg2/tx2): start from the master's
        # clrMap, then apply any override on the layout or the slide —
        # dark designs swap text/background roles via clrMapOvr, and
        # ignoring the override renders light text as dark.
        alias_map: dict = {}
        clr_map = master.element.find(f"{{{_P_NS}}}clrMap")
        if clr_map is not None:
            alias_map.update(clr_map.attrib)
        for owner in (src_slide.slide_layout, src_slide):
            ovr = owner._element.find(
                f"{{{_P_NS}}}clrMapOvr/{{{_A_NS}}}overrideClrMapping"
            )
            if ovr is not None:
                alias_map.update(ovr.attrib)
        for alias, target in alias_map.items():
            if target in colors:
                colors[alias] = colors[target]

        fonts: dict = {}
        font_scheme = theme_root.find(f".//{{{_A_NS}}}fontScheme")
        for scope, tag in (("mj", "majorFont"), ("mn", "minorFont")):
            group = (
                font_scheme.find(f"{{{_A_NS}}}{tag}")
                if font_scheme is not None else None
            )
            if group is None:
                continue
            for kind in ("latin", "ea", "cs"):
                el = group.find(f"{{{_A_NS}}}{kind}")
                if el is not None and el.get("typeface"):
                    suffix = "lt" if kind == "latin" else kind
                    fonts[f"+{scope}-{suffix}"] = el.get("typeface")

        return colors, fonts
    except Exception:
        return {}, {}


def _bake_theme_references(slide_element, colors: dict, fonts: dict) -> None:
    """Replace scheme colors / theme fonts with their literal values.

    ``<a:schemeClr val="accent1">`` becomes ``<a:srgbClr val="AABBCC">``
    (children — lumMod, alpha, etc. — are preserved, and srgbClr
    accepts the same transform children). Theme-font tokens on any
    ``typeface`` attribute are swapped for the resolved face. Unknown
    references are left untouched.
    """
    if colors:
        for sc in slide_element.iter(f"{{{_A_NS}}}schemeClr"):
            hex_val = colors.get(sc.get("val"))
            if hex_val:
                sc.tag = f"{{{_A_NS}}}srgbClr"
                sc.set("val", hex_val)
    if fonts:
        for node in slide_element.iter():
            face = fonts.get(node.get("typeface") or "")
            if face:
                node.set("typeface", face)


def _bake_placeholder_text_colors(slide_element, src_slide, colors: dict) -> None:
    """Give colorless placeholder runs their inherited color, explicitly."""
    from copy import deepcopy
    from lxml import etree

    for sp in slide_element.iter(f"{{{_P_NS}}}sp"):
        ph = sp.find(f".//{{{_P_NS}}}ph")
        if ph is None:
            continue
        hex_color = _inherited_placeholder_color(src_slide, ph, colors)
        if not hex_color:
            continue
        fill = etree.SubElement(
            etree.Element("dummy"), f"{{{_A_NS}}}solidFill"
        )
        clr = etree.SubElement(fill, f"{{{_A_NS}}}srgbClr")
        clr.set("val", hex_color)

        for para in sp.iter(f"{{{_A_NS}}}p"):
            for run in para.findall(f"{{{_A_NS}}}r"):
                rpr = run.find(f"{{{_A_NS}}}rPr")
                if rpr is None:
                    rpr = etree.Element(f"{{{_A_NS}}}rPr")
                    run.insert(0, rpr)
                if rpr.find(f"{{{_A_NS}}}solidFill") is None and \
                   rpr.find(f"{{{_A_NS}}}gradFill") is None:
                    ln = rpr.find(f"{{{_A_NS}}}ln")
                    pos = list(rpr).index(ln) + 1 if ln is not None else 0
                    rpr.insert(pos, deepcopy(fill))


def _inherited_placeholder_color(src_slide, ph, colors: dict):
    """Resolve a placeholder's inherited text color to a hex string.

    Checks the source layout's matching placeholder list-style first,
    then the master's txStyles (title/body/other by placeholder type).
    Scheme references resolve through the already-override-aware color
    map. Returns None when nothing explicit is found.
    """
    ph_type = ph.get("type") or "body"
    ph_idx = ph.get("idx")

    def fill_hex(def_rpr):
        if def_rpr is None:
            return None
        solid = def_rpr.find(f"{{{_A_NS}}}solidFill")
        if solid is None:
            return None
        srgb = solid.find(f"{{{_A_NS}}}srgbClr")
        if srgb is not None:
            return srgb.get("val")
        scheme = solid.find(f"{{{_A_NS}}}schemeClr")
        if scheme is not None:
            return colors.get(scheme.get("val"))
        return None

    def matches(candidate_ph) -> bool:
        c_type = candidate_ph.get("type") or "body"
        if ph_type in ("title", "ctrTitle"):
            return c_type in ("title", "ctrTitle")
        return c_type == ph_type and candidate_ph.get("idx") == ph_idx

    try:
        # 1. Placeholder SHAPES, in inheritance order: the layout's,
        # then the master's. Google Slides puts the effective color on
        # these shapes' list styles far more often than on txStyles.
        layout = src_slide.slide_layout
        for owner in (layout, layout.slide_master):
            for sp in owner._element.iter(f"{{{_P_NS}}}sp"):
                cph = sp.find(f".//{{{_P_NS}}}ph")
                if cph is None or not matches(cph):
                    continue
                lst = sp.find(f".//{{{_A_NS}}}lstStyle")
                if lst is not None:
                    for lvl in lst:
                        hex_color = fill_hex(lvl.find(f"{{{_A_NS}}}defRPr"))
                        if hex_color:
                            return hex_color

        # 2. Master txStyles by placeholder role.
        style_tag = (
            "titleStyle" if ph_type in ("title", "ctrTitle")
            else "bodyStyle" if ph_type in ("body", "subTitle")
            else "otherStyle"
        )
        tx = layout.slide_master.element.find(f"{{{_P_NS}}}txStyles")
        style = tx.find(f"{{{_P_NS}}}{style_tag}") if tx is not None else None
        if style is not None:
            lvl1 = style.find(f"{{{_A_NS}}}lvl1pPr")
            if lvl1 is not None:
                return fill_hex(lvl1.find(f"{{{_A_NS}}}defRPr"))
    except Exception:
        return None
    return None
