"""Agent-built slides: an autonomous Claude Code session (Fable, via
claude-agent-sdk) authors a one-slide .pptx with native, Google-Slides-
editable shapes.

Two styles share the engine:
- "open": the spec's description gives the agent creative authority over the
  layout of already-resolved content (images matched, text bound upstream).
- "adaptive": the agent starts from a COPY of a reference skeleton — swapping
  the provided content into it — and may adjust geometry only where the
  actual content misfits; fonts, colors, and decorations must not change.

Autonomy inside, guarantees outside: the agent gets Write/Edit/Bash/Read in
an isolated workspace and is told to render (LibreOffice) and LOOK at its
own output; the result must then pass a deterministic acceptance gate
(opens, one slide, every required text verbatim, images embedded, no
leftover {{tokens}}) plus a vision check. One retry with the findings; a
second failure returns None and the caller falls back to the deterministic
freeform layout so the deck always completes.
"""

from __future__ import annotations

import json
import shutil
import subprocess
from pathlib import Path

from core import llm
from deck import soffice
from deck.assemble import MSO_PICTURE, TOKEN_RE, iter_shapes
from deck.brand import BRAND_DIR, lint_presentation, load_brand, reference_renders

MAX_TURNS = 50
# Batch sessions read a render per slide before they can even start judging,
# so they need materially more room than a single-slide session.
BATCH_MAX_TURNS = 80
ATTEMPTS = 2
# A slide session writes, renders and iterates; give it room but never hang.
CLI_SESSION_TIMEOUT_S = 1800


def agent_model() -> str:
    """Model for autonomous slide/design sessions. Fable by default;
    SG_AGENT_MODEL overrides it (e.g. when a subscription's Fable quota is
    exhausted — subscription limits are per-model)."""
    import os

    return os.environ.get("SG_AGENT_MODEL") or llm.FABLE

BATCH_PROMPT = """You are building {n} presentation slides for ONE deck, in this directory.
Work autonomously until they are genuinely good.

`slides.json` lists every slide to build, in deck order. For each entry write
`out/<id>.pptx` — a valid ONE-slide PowerPoint file.

Hard rules (per slide):
- NATIVE shapes only (text boxes, pictures, simple shapes) so the slides stay
  fully editable in Google Slides. Never rasterise text into an image.
- Slide size {width}in x {height}in.
- Every value in that slide's `texts` must appear VERBATIM — never altered,
  truncated or paraphrased. Where `step_no` is present the slide belongs to the
  deck's numbered run; title it exactly the way `design_guide.md` says numbered
  slides are titled.
- Every file in that slide's `images` must be placed, aspect ratio preserved
  (fit inside its frame; never stretch or distort).
- Leave no `{{{{token}}}}` placeholder text anywhere.

CONSISTENCY IS THE POINT — read this twice:
- These slides sit near each other in one deck. Slides of the same family MUST
  share their geometry and typography EXACTLY, differing only in content.
- So write ONE parameterised builder function and call it per slide. Do not
  hand-position each slide separately; that is how drift happens.
- `design_guide.md` states this deck's conventions, measured from the real
  templates. FOLLOW IT — it outranks your own layout instincts.
- `context/` holds renders of the FIXED slides that sit immediately before and
  after these in the deck (each slide's entry names its own). Read them FIRST:
  your slides must look like they belong between them. If one of those layouts
  fits your content, reproduce its geometry rather than inventing your own.
- `brand/` holds the brand kit, reference renders and the real logo files. Set
  brand fonts and colours EXPLICITLY on every run; place logos from
  `brand/logos/` and never draw one.

Method:
- Write one build script; run it with:
  `uv run --with python-pptx --with pillow python build.py`
- Then RENDER AND LOOK. Convert each slide, e.g.
  `soffice {profile_flag} --headless --convert-to png out/<id>.pptx --outdir out/`
  (the profile flag is required — without it concurrent LibreOffice runs fail),
  then Read the PNGs. Judge two things: each slide on its own (overlaps,
  crowding, unreadable sizes, stretched images) AND the slides TOGETHER — same
  family must be pixel-consistent in layout. Fix the script and re-render until
  both hold.
- Finish only when every `out/<id>.pptx` exists and you are satisfied.
{feedback}"""

PROMPT_COMMON = """You are building ONE presentation slide as a file named `slide.pptx` in the
current directory. Work autonomously until it is genuinely good.

Hard rules:
- `slide.pptx` must be a valid one-slide PowerPoint file made of NATIVE shapes
  (text boxes, pictures, simple shapes) so it stays fully editable in Google
  Slides. Never rasterize text into images.
- Slide size: {width}in x {height}in.
- `content.json` in this directory holds the content contract: every text
  value in it must appear on the slide VERBATIM (you may add small structural
  labels of your own, but never alter, truncate, or paraphrase the provided
  values), and every listed image file must be placed, aspect ratio
  preserved (fit inside its frame; never stretch).
- If any template token like {{{{name}}}} would remain with no provided value,
  remove that token text cleanly.
- Brand: the rules are in content.json under "brand", and the `brand/`
  directory holds the ground truth — BEFORE designing, Read the reference
  renders in `brand/reference_slides/` and match their visual family
  (palette, type hierarchy, whitespace). When a logo belongs on the slide,
  place a file from `brand/logos/` — never draw or fabricate one. Set the
  brand fonts and colors EXPLICITLY on every text run.
- Deck conventions: `design_guide.md` (when present) states how this deck's
  slide families are built — title wording and treatment, body copy
  placement and colour, image proportions, numbering. FOLLOW IT; it outranks
  your own layout instincts. Where content.json provides a step number, this
  slide belongs to the deck's numbered run — title it exactly the way the
  guide says numbered slides are titled.
- Neighbours: when `neighbours/` exists, Read the PNG render(s) there FIRST.
  Those are the slides immediately before and after this one; your slide has
  to look like it belongs between them. If a neighbour's template
  (`neighbours/*_template.pptx`) is a good fit for your content, START FROM
  THAT FILE — load it, swap in the provided content, adjust only what
  genuinely misfits — rather than composing a layout from scratch. Reusing
  the deck's own layout is always better than re-deriving one.

Method:
- Write a Python build script and run it with:
  `uv run --with python-pptx --with pillow python build.py`
- Then RENDER and LOOK at your work:
  `soffice --headless --convert-to png slide.pptx` produces `slide.png` —
  Read it and judge honestly: overlaps, crowding, unreadable sizes, awkward
  whitespace, misaligned elements. Iterate script -> render -> look until it
  is clean. If `soffice` is unavailable, be extra conservative with spacing.
- Finish only when `slide.pptx` exists and you are satisfied with the render.
"""

PROMPT_OPEN = """
Layout brief (you have creative authority over arrangement and styling within
the brand kit):
{description}
"""

PROMPT_ADAPTIVE = """
This slide ADAPTS a reference layout. `reference.pptx` is the reference slide
and `reference_layout.json` inventories its shapes (geometry in inches, fonts,
which shapes are content vs decoration).

Method for this style — follow it strictly:
- START FROM THE REFERENCE: load `reference.pptx`, replace the content in its
  text/token shapes and image slots with the provided content, and save as
  `slide.pptx`. Decorations, background, fonts, and colors carry over
  automatically — do not rebuild them and do not restyle them.
- You MAY adjust geometry (grow/shift/shrink boxes, rebalance columns) ONLY
  where the actual content misfits the original frame — text overflowing,
  an image aspect fighting its frame. Keep every adjustment minimal; when in
  doubt, copy the reference exactly.
{description_block}"""

VERIFY_SCHEMA = {
    "type": "object",
    "properties": {"match": {"type": "boolean"}, "reason": {"type": "string"}},
    "required": ["match", "reason"],
    "additionalProperties": False,
}

VERIFY_PROMPT = """The FIRST {n_ref} image(s) are reference slides from the canonical brand
deck. The LAST image is a render of a generated slide for the same
customer-facing test report.

It was built to this brief: {brief}

IMPORTANT — the small numeral in the top-left corner (inside or beside the
purple rail) is PowerPoint's slide-number placeholder, inherited from the
template. Its value is assigned by position when the full deck is assembled.
You are looking at this slide rendered ALONE, so it always reads "1" no
matter where the slide will actually sit. Do NOT treat that numeral, or its
disagreement with a "Step N" title, as a defect — it is an artifact of this
preview and cannot occur in the delivered deck. (If the slide has drawn its
OWN extra step number somewhere, that duplicate IS a defect — report it.)

Judge the generated slide for:
1. Defects a customer would notice: overlapping elements, text cut off or
   unreadably small, images badly squeezed or overflowing the slide, or the
   slide plainly not delivering the brief.
2. Brand family: it must plausibly belong to the same deck as the reference
   slides — palette, typographic hierarchy, whitespace, logo usage.
Stylistic taste differences are NOT defects. Answer match=false only for
real defects or clear brand breaks, with the reason."""


def _shape_inventory(pptx_path: Path) -> list[dict]:
    from pptx import Presentation
    from pptx.util import Emu

    out = []
    slide = Presentation(str(pptx_path)).slides[0]
    for shape in iter_shapes(slide):
        text = shape.text_frame.text if shape.has_text_frame else ""
        entry = {
            "name": shape.name,
            "kind": "picture" if shape.shape_type == MSO_PICTURE else "shape",
            "x_in": round(Emu(shape.left).inches, 2),
            "y_in": round(Emu(shape.top).inches, 2),
            "w_in": round(Emu(shape.width).inches, 2),
            "h_in": round(Emu(shape.height).inches, 2),
            "text": text[:200],
            "is_content": bool(TOKEN_RE.search(text))
            or shape.shape_type == MSO_PICTURE
            or "insert screenshot" in text.lower(),
        }
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs[:1]:
                    entry["font"] = {
                        "name": run.font.name,
                        "size_pt": run.font.size.pt if run.font.size else None,
                        "bold": run.font.bold,
                    }
                    break
                break
        out.append(entry)
    return out


def prepare_workspace(
    workspace: Path,
    style: str,
    description: str,
    texts: dict[str, str],
    images: list[str],
    skeleton: str | None,
    brief: dict | None = None,
) -> None:
    workspace.mkdir(parents=True, exist_ok=True)
    # Deck conventions (prose) + where this slide sits (facts + neighbour
    # renders) — cohesion comes from seeing the slides it will sit between.
    from deck.design import load_design_guide, write_brief

    write_brief(workspace, load_design_guide(), brief or {})
    # Brand assets travel with the workspace: rules + reference renders the
    # agent can LOOK at + the real logo files (never fabricate a logo).
    shutil.copytree(BRAND_DIR, workspace / "brand", dirs_exist_ok=True)
    from PIL import Image

    image_entries = []
    for i, src in enumerate(images):
        dest = workspace / f"image_{i}{Path(src).suffix.lower()}"
        shutil.copy(src, dest)
        with Image.open(dest) as im:
            w, h = im.size
        image_entries.append({"file": dest.name, "width_px": w, "height_px": h})
    brand = load_brand()
    content = {
        "style": style,
        "description": description,
        "texts": texts,
        "images": image_entries,
        "brand": brand,
        "slide_size_in": brand["slide_size_in"],
    }
    (workspace / "content.json").write_text(json.dumps(content, indent=2))
    if style == "adaptive":
        if not skeleton:
            raise ValueError("adaptive agent slide requires a reference skeleton")
        shutil.copy(skeleton, workspace / "reference.pptx")
        (workspace / "reference_layout.json").write_text(
            json.dumps(_shape_inventory(Path(skeleton)), indent=2)
        )


def _prompt(style: str, description: str) -> str:
    w, h = load_brand()["slide_size_in"]
    common = PROMPT_COMMON.format(width=w, height=h)
    if style == "adaptive":
        block = (
            f"\nAdditional guidance:\n{description}\n" if description.strip() else "\n"
        )
        return common + PROMPT_ADAPTIVE.format(description_block=block)
    return common + PROMPT_OPEN.format(description=description)


def run_agent_session(
    workspace: Path, prompt: str, log=print, max_turns: int = MAX_TURNS
) -> tuple[int, str | None]:
    """One autonomous Claude Code session in an isolated workspace, on the
    preferred agent model — falling down the tier ladder if it is
    unavailable. Returns (turns, error). Shared by slide building and
    design-guide derivation."""
    requested = agent_model()
    chain = llm.fallback_chain(requested)
    last_error = None
    for i, candidate in enumerate(chain):
        turns, error = _session_once(
            workspace, prompt, log, max_turns, candidate, transport=session_transport()
        )
        if error and llm.is_availability_issue(error) and i < len(chain) - 1:
            last_error = error
            continue
        if candidate != requested and not error:
            llm.record_substitution(requested, candidate, last_error or "unavailable", log)
        return turns, error
    return 0, last_error


def session_transport() -> str:
    """How to run an agent session: "cli" shells out to the installed
    `claude` binary in headless mode; "sdk" uses claude-agent-sdk's managed
    process. Both do the same job — write, run, render, iterate — but they
    are metered differently, so the transport follows the selected LLM
    backend (and SG_AGENT_TRANSPORT overrides for debugging)."""
    import os

    override = os.environ.get("SG_AGENT_TRANSPORT")
    if override in ("cli", "sdk"):
        return override
    return "sdk" if llm.backend().name == "agent-sdk" else "cli"


def _session_once(
    workspace: Path,
    prompt: str,
    log,
    max_turns: int,
    model: str,
    transport: str = "sdk",
) -> tuple[int, str | None]:
    import os

    # Safety net for test suites: an "offline" test that forgets to stub the
    # session would otherwise spawn real agents and burn subscription quota.
    if os.environ.get("SG_AGENT_OFFLINE") == "1":
        raise AssertionError(
            "a real agent session was attempted while SG_AGENT_OFFLINE=1 — "
            "stub deck.agent_slide._session_once in this test"
        )
    if transport == "cli":
        return _session_once_cli(workspace, prompt, log, max_turns, model)
    return _session_once_sdk(workspace, prompt, log, max_turns, model)


def _session_once_cli(
    workspace: Path, prompt: str, log, max_turns: int, model: str
) -> tuple[int, str | None]:
    """One headless `claude -p` session with tools, in the workspace."""
    import json as _json
    import shutil
    import subprocess

    from core.llm import is_availability_issue

    exe = shutil.which("claude")
    if not exe:
        return 0, "the `claude` CLI was not found on PATH"
    cmd = [
        exe, "-p",
        "--model", model,
        "--tools", "Read,Write,Edit,Bash",
        "--permission-mode", "bypassPermissions",
        "--max-turns", str(max_turns),
        "--output-format", "json",
    ]
    try:
        proc = subprocess.run(
            cmd, input=prompt, capture_output=True, text=True,
            cwd=str(workspace), timeout=CLI_SESSION_TIMEOUT_S,
        )
    except subprocess.TimeoutExpired:
        return 0, f"session exceeded {CLI_SESSION_TIMEOUT_S}s"
    if proc.returncode != 0 and not proc.stdout.strip():
        detail = (proc.stderr or "").strip()[:300]
        return 0, f"claude CLI failed (exit {proc.returncode}): {detail}"
    try:
        payload = _json.loads(proc.stdout)
    except _json.JSONDecodeError:
        return 0, f"claude CLI returned non-JSON output: {proc.stdout[:200]!r}"
    turns = int(payload.get("num_turns") or 0)
    if not payload.get("is_error"):
        return turns, None
    detail = str(payload.get("result") or "session reported an error")[:300]
    if is_availability_issue(detail):
        log(f"  {detail}")
    return turns, detail


def _session_once_sdk(
    workspace: Path, prompt: str, log, max_turns: int, model: str
) -> tuple[int, str | None]:
    from claude_agent_sdk import (
        AssistantMessage,
        ClaudeAgentOptions,
        RateLimitEvent,
        ResultMessage,
        query,
    )

    from core.llm import SDK_BUFFER_BYTES, rate_limit_note, run_coro_in_thread

    options = ClaudeAgentOptions(
        model=model,
        cwd=str(workspace),
        tools=["Read", "Write", "Edit", "Bash"],
        setting_sources=[],
        permission_mode="bypassPermissions",
        max_turns=max_turns,
        # These sessions Read full-page PNG renders to compare slides against
        # each other; the SDK's 1MB stdio default would fatally kill the
        # session mid-batch on a large render.
        max_buffer_size=SDK_BUFFER_BYTES,
    )
    state = {"turns": 0, "error": None}

    async def run():
        async for msg in query(prompt=prompt, options=options):
            if isinstance(msg, AssistantMessage):
                state["turns"] += 1
            if isinstance(msg, RateLimitEvent):
                note = rate_limit_note(msg)
                if note and not state["error"]:
                    # Record and keep draining: returning mid-iteration closes
                    # the SDK's async generator while it is still running.
                    state["error"] = note
                    log(f"  {note}")
            if isinstance(msg, ResultMessage) and msg.is_error:
                state["error"] = state["error"] or str(msg.result)[:300]

    try:
        run_coro_in_thread(run)
    except Exception as e:
        # A recorded rate-limit note is more actionable than the SDK's
        # downstream exception, so it wins.
        state["error"] = state["error"] or f"{type(e).__name__}: {e}"
    return state["turns"], state["error"]


def _run_agent(workspace: Path, prompt: str, log) -> tuple[int, str | None]:
    """Slide-building session (test seam: patched in offline tests)."""
    return run_agent_session(workspace, prompt, log=log)


def gate(pptx_path: Path, required_texts: list[str], min_images: int) -> list[str]:
    """Deterministic acceptance checks. Empty list = pass."""
    issues: list[str] = []
    if not pptx_path.exists():
        return ["slide.pptx was not produced"]
    try:
        from pptx import Presentation

        pres = Presentation(str(pptx_path))
    except Exception as e:
        return [f"slide.pptx does not open as a presentation: {e}"]
    if len(pres.slides) != 1:
        issues.append(f"expected exactly 1 slide, found {len(pres.slides)}")
    if not pres.slides:
        return issues
    slide = pres.slides[0]
    norm = lambda s: " ".join(s.split())
    slide_text = norm(
        " ".join(
            sh.text_frame.text for sh in iter_shapes(slide) if sh.has_text_frame
        )
    )
    for req in required_texts:
        if norm(req) and norm(req) not in slide_text:
            issues.append(f"required text missing verbatim: {req[:80]!r}")
    if "{{" in slide_text:
        issues.append("unresolved {{token}} text remains on the slide")
    pics = [sh for sh in iter_shapes(slide) if sh.shape_type == MSO_PICTURE]
    if len(pics) < min_images:
        issues.append(f"expected >= {min_images} embedded image(s), found {len(pics)}")
    if not issues:
        for f in lint_presentation(pptx_path):
            issues.append(f"brand lint [{f['check']}]: {f['detail']}")
    return issues


def _render_png(pptx_path: Path) -> bytes | None:
    """Render for the acceptance verdict. FIDELITY: this render decides
    whether a slide ships, so it should show the fonts and text metrics the
    engineer will actually see, not LibreOffice's substitutes."""
    from deck import render

    produced = render.convert(
        pptx_path, "png", pptx_path.parent, purpose=render.FIDELITY,
        log=lambda *a: None,
    )
    return produced.read_bytes() if produced else None


def _vision_verdict(pptx_path: Path, brief: str) -> dict:
    png = _render_png(pptx_path)
    if png is None:
        return {"match": True, "reason": "no renderer available; gate-only acceptance"}
    refs = [p.read_bytes() for p in reference_renders()[:2]]
    try:
        return llm.complete(
            VERIFY_PROMPT.format(
                n_ref=len(refs), brief=" ".join(brief.split())[:600]
            ),
            schema=VERIFY_SCHEMA,
            images=refs + [png],
            max_tokens=1000,
            model=agent_model(),
        )
    except Exception as e:
        return {"match": True, "reason": f"vision check unavailable ({e}); gate-only"}


def build_agent_slides(
    jobs: list[dict],
    work_root: Path,
    group_brief: dict | None = None,
    log=print,
) -> dict[str, dict]:
    """Build every agent slide of a run in ONE session, then gate each
    independently. Returns {sid: report}.

    One session rather than N: these slides must be consistent with each
    other, and only an agent that can see them all together can check that.
    It also amortises the guide/brand/context reading and the build script
    across the whole group, which makes it faster than N sessions as well.

    Failure is per slide, not per batch: slides that pass the gate are kept,
    a single retry re-prompts naming only the failures, and whatever still
    fails is reported so the caller can fall back individually.

    NOTE for the future: at ~30+ agent slides a single session becomes
    unwieldy (turn budget, context). The natural extension is batching BY
    FAMILY — all numbered steps in one session, all stat cards in another —
    which preserves the consistency benefit where it actually matters. Not
    needed at current deck sizes.
    """
    if not jobs:
        return {}
    workspace = work_root / "batch"
    prepare_batch_workspace(workspace, jobs, group_brief or {})
    out_dir = workspace / "out"
    # `session_turns` is the shared cost of the batch, recorded once per
    # attempt — not summed into every slide, which would inflate it N-fold.
    reports = {
        j["sid"]: {"pptx": None, "attempts": 0, "session_turns": [], "issues": [],
                   "workspace": str(workspace)}
        for j in jobs
    }
    by_sid = {j["sid"]: j for j in jobs}
    pending = [j["sid"] for j in jobs]
    feedback = ""

    for attempt in range(1, ATTEMPTS + 1):
        log(
            f"  agent slides: attempt {attempt} for {len(pending)} slide(s) "
            f"(model {agent_model()}, one session)"
        )
        w, h = load_brand()["slide_size_in"]
        prompt = BATCH_PROMPT.format(
            n=len(pending), width=w, height=h,
            profile_flag=soffice.PROFILE_FLAG_HINT, feedback=feedback,
        )
        # Through run_agent_session, so the batch inherits the model-tier
        # fallback ladder rather than dying on one unavailable tier.
        turns, error = run_agent_session(
            workspace, prompt, log=log, max_turns=BATCH_MAX_TURNS
        )
        still_failing: list[str] = []
        for sid in pending:
            job = by_sid[sid]
            reports[sid]["attempts"] = attempt
            # One entry per attempt this slide took part in.
            reports[sid]["session_turns"].append(turns)
            path = out_dir / f"{sid}.pptx"
            required = [
                v for v in job["texts"].values() if v and v.strip() and v.strip() != "—"
            ]
            issues = gate(path, required, min_images=len(job["images"]))
            if issues and error and not path.exists():
                issues.insert(0, f"session error: {error}")
            if issues:
                reports[sid]["issues"] = issues
                still_failing.append(sid)
                continue
            verdict = _vision_verdict(path, job.get("description") or job["sid"])
            reports[sid]["vision"] = verdict
            if verdict["match"]:
                reports[sid]["pptx"] = str(path)
                log(f"    accepted {sid}")
            else:
                reports[sid]["issues"] = [
                    f"vision review found a defect: {verdict['reason']}"
                ]
                still_failing.append(sid)
        if not still_failing:
            return reports
        log(
            f"    {len(still_failing)} slide(s) rejected: "
            + ", ".join(still_failing)
        )
        pending = still_failing
        feedback = "\nYour previous attempt was rejected for these slides ONLY — "
        feedback += "leave the accepted ones untouched and fix these:\n"
        for sid in pending:
            feedback += f"- {sid}: {'; '.join(reports[sid]['issues'])[:300]}\n"
    return reports


def _slide_sized(src: str) -> str:
    """Downscale to what a full-width slide can actually show. The widest
    image area in these layouts is ~8.8in, so EMBED_DPI across that is the
    practical ceiling for anything placed on a slide."""
    from pptx.util import Inches

    from deck.assemble import sized_for_slot

    return sized_for_slot(src, Inches(8.8), Inches(5.0))


def prepare_batch_workspace(
    workspace: Path, jobs: list[dict], group_brief: dict
) -> None:
    """One workspace for the whole group: brand kit and design guide once,
    deduped context renders once, and slides.json describing every slide."""
    from PIL import Image

    from deck.design import load_design_guide

    workspace.mkdir(parents=True, exist_ok=True)
    # Start from an empty out/: a slide left behind by an earlier run would
    # otherwise pass the gate and ship as if it had just been built.
    out = workspace / "out"
    if out.exists():
        shutil.rmtree(out)
    out.mkdir()
    shutil.copytree(BRAND_DIR, workspace / "brand", dirs_exist_ok=True)
    guide = load_design_guide()
    if guide:
        (workspace / "design_guide.md").write_text(guide)

    context = workspace / "context"
    copied: dict[str, str] = {}
    for render in group_brief.get("renders", []):
        src = Path(render)
        if not src.exists() or render in copied:
            continue
        context.mkdir(exist_ok=True)
        shutil.copy(src, context / src.name)
        copied[render] = f"context/{src.name}"

    entries = []
    for job in jobs:
        sid = job["sid"]
        slide_dir = workspace / "assets" / sid
        slide_dir.mkdir(parents=True, exist_ok=True)
        images = []
        for i, src in enumerate(job["images"]):
            dest = slide_dir / f"image_{i}{Path(src).suffix.lower()}"
            # Slide-sized copy, not the native capture: the agent both reads
            # and embeds this file, so a 3840px original would bloat the
            # produced slide and slow every read of it.
            shutil.copy(_slide_sized(src), dest)
            with Image.open(dest) as im:
                w, h = im.size
            images.append(
                {"file": f"assets/{sid}/{dest.name}", "width_px": w, "height_px": h}
            )
        info = (group_brief.get("per_slide") or {}).get(sid) or {}
        neighbours = {}
        for side in ("previous", "next"):
            n = info.get(side)
            if not n:
                continue
            neighbours[side] = {
                "id": n.get("id"),
                "title": n.get("title"),
                "step_no": n.get("step_no"),
                "render": copied.get(n.get("render") or ""),
            }
        entries.append(
            {
                "id": sid,
                "output": f"out/{sid}.pptx",
                "style": job.get("style", "open"),
                "brief": job.get("description", ""),
                "step_no": job["texts"].get("step_no"),
                "texts": job["texts"],
                "images": images,
                "neighbours": neighbours,
            }
        )
    (workspace / "slides.json").write_text(
        json.dumps({"brand": load_brand(), "slides": entries}, indent=2)
    )


def build_agent_slide(
    sid: str,
    style: str,
    description: str,
    texts: dict[str, str],
    images: list[str],
    skeleton: str | None,
    work_root: Path,
    log=print,
    brief: dict | None = None,
) -> dict:
    """Returns {"pptx": str|None, "attempts", "turns", "issues", "workspace"}."""
    workspace = work_root / sid
    prepare_workspace(workspace, style, description, texts, images, skeleton, brief)
    slide_path = workspace / "slide.pptx"
    required = [v for v in texts.values() if v and v.strip() and v.strip() != "—"]
    report = {"pptx": None, "attempts": 0, "turns": 0, "issues": [], "workspace": str(workspace)}

    prompt = _prompt(style, description)
    for attempt in range(1, ATTEMPTS + 1):
        report["attempts"] = attempt
        log(f"  agent slide {sid} ({style}): attempt {attempt} (model {agent_model()})")
        turns, error = _run_agent(workspace, prompt, log)
        report["turns"] += turns
        issues = gate(slide_path, required, min_images=len(images))
        if error and not slide_path.exists():
            issues.insert(0, f"agent session error: {error}")
        if not issues:
            verdict = _vision_verdict(slide_path, description or style)
            report["vision"] = verdict
            if verdict["match"]:
                report["pptx"] = str(slide_path)
                log(f"  agent slide {sid}: accepted after {report['turns']} turns")
                return report
            issues = [f"vision review found a defect: {verdict['reason']}"]
        report["issues"] = issues
        log(f"  agent slide {sid}: attempt {attempt} rejected: {'; '.join(issues)[:150]}")
        prompt = (
            _prompt(style, description)
            + "\n\nYour previous attempt (files still in this directory) was rejected:\n- "
            + "\n- ".join(issues)
            + "\nFix these specific problems and finish again."
        )
    return report
