"""Brand kit loading and enforcement.

The kit lives in deck/brand/ (brand.yaml + reference_slides/ renders +
logos/) and is the single source of truth for every producer and checker:
the freeform renderer's defaults, the agent-slide prompts and workspace
assets, the copywriter's voice rules, the per-slide acceptance gate, and
the post-assembly deck audit.

Two enforcement tiers:
- lint_presentation(): deterministic, token-free checks over actual pptx
  shapes — explicit fonts outside the allowed set, explicit colors outside
  the palette, sub-minimum font sizes, badly out-of-bounds content,
  stretched images, likely text overflow. Runs on every slide of every deck.
- review_render(): a vision check of a rendered slide, anchored on the
  reference renders ("same visual family?"). Costs a model call; used in
  the agent-slide gate and (optionally) the deck audit.

audit_deck() combines them into brand_report.json — report-only: assembled
skeleton slides can't be regenerated, so findings inform rather than block.
"""

from __future__ import annotations

import json
from functools import lru_cache
from pathlib import Path

import yaml

from core import llm

BRAND_DIR = Path(__file__).resolve().parent / "brand"

# Concurrency for the audit's per-slide vision reviews (independent calls;
# modest to stay friendly to subscription rate limits).
AUDIT_WORKERS = 4


@lru_cache(maxsize=1)
def load_brand() -> dict:
    return yaml.safe_load((BRAND_DIR / "brand.yaml").read_text())


def reference_renders() -> list[Path]:
    return sorted((BRAND_DIR / "reference_slides").glob("*.png"))


def _hex_to_rgb(value: str) -> tuple[int, int, int]:
    v = value.lstrip("#")
    return tuple(int(v[i : i + 2], 16) for i in (0, 2, 4))


def _color_ok(rgb, brand: dict) -> bool:
    tol = brand["colors"].get("tolerance", 40)
    candidates = [
        _hex_to_rgb(v) for k, v in brand["colors"].items() if isinstance(v, str)
    ]
    # Neutrals (greys) are always acceptable chrome.
    r, g, b = rgb
    if max(r, g, b) - min(r, g, b) <= 24:
        return True
    return any(
        max(abs(r - cr), abs(g - cg), abs(b - cb)) <= tol for cr, cg, cb in candidates
    )


def lint_presentation(pptx_path: Path, brand: dict | None = None) -> list[dict]:
    """Findings: [{slide (1-based), check, detail}]. Empty = clean."""
    from pptx import Presentation
    from pptx.util import Emu

    from deck.assemble import MSO_PICTURE, iter_shapes

    brand = brand or load_brand()
    allowed_fonts = set(brand["fonts"]["allowed"])
    min_pt = brand["typography"]["min_pt"]
    overhang = brand["layout"]["overhang_tolerance_in"]
    max_distort = brand["layout"]["max_image_aspect_distortion"]
    findings: list[dict] = []
    pres = Presentation(str(pptx_path))
    slide_w, slide_h = Emu(pres.slide_width).inches, Emu(pres.slide_height).inches

    for idx, slide in enumerate(pres.slides, start=1):
        def flag(check: str, detail: str) -> None:
            findings.append({"slide": idx, "check": check, "detail": detail})

        for shape in iter_shapes(slide):
            if None in (shape.left, shape.top, shape.width, shape.height):
                continue
            x, y = Emu(shape.left).inches, Emu(shape.top).inches
            w, h = Emu(shape.width).inches, Emu(shape.height).inches
            is_full_bleed = w >= slide_w * 0.95 or h >= slide_h * 0.95
            has_text = shape.has_text_frame and shape.text_frame.text.strip()
            if (has_text or shape.shape_type == MSO_PICTURE) and not is_full_bleed:
                if (
                    x < -overhang
                    or y < -overhang
                    or x + w > slide_w + overhang
                    or y + h > slide_h + overhang
                ):
                    flag(
                        "out_of_bounds",
                        f"{'text' if has_text else 'image'} shape at "
                        f"({x:.1f},{y:.1f}) {w:.1f}x{h:.1f}in exceeds the slide",
                    )
            if shape.shape_type == MSO_PICTURE:
                crop = (
                    getattr(shape, "crop_left", 0)
                    or getattr(shape, "crop_right", 0)
                    or getattr(shape, "crop_top", 0)
                    or getattr(shape, "crop_bottom", 0)
                )
                try:
                    nat_w, nat_h = shape.image.size
                except Exception:
                    nat_w = nat_h = 0
                if not crop and nat_w and nat_h and h > 0 and nat_h > 0:
                    distortion = abs((w / h) / (nat_w / nat_h) - 1)
                    if distortion > max_distort:
                        flag(
                            "stretched_image",
                            f"picture aspect off by {distortion:.0%} vs its source",
                        )
                continue
            if not shape.has_text_frame:
                continue
            text_chars = 0
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    text_chars += len(run.text)
                    name = run.font.name
                    if name and name not in allowed_fonts and run.text.strip():
                        flag("font", f"non-brand font {name!r}: {run.text.strip()[:40]!r}")
                    size = run.font.size
                    # Micro-strings (page numbers, decorative numerals) are
                    # legitimate at tiny sizes in the canonical templates.
                    if size and size.pt < min_pt and len(run.text.strip()) > 3:
                        flag(
                            "min_font_size",
                            f"{size.pt:.0f}pt below minimum {min_pt}pt: "
                            f"{run.text.strip()[:40]!r}",
                        )
                    try:
                        rgb = run.font.color.rgb
                    except Exception:
                        rgb = None
                    if rgb is not None and run.text.strip():
                        triplet = (rgb[0], rgb[1], rgb[2])
                        if not _color_ok(triplet, brand):
                            flag(
                                "color",
                                f"off-palette text color #{str(rgb)}: "
                                f"{run.text.strip()[:40]!r}",
                            )
            # Overflow heuristic: conservative — flags only clear offenders.
            # Small declared heights are auto-grow/auto-fit frames in the
            # Google-exported templates; their box size is not the render size.
            if text_chars and w > 0 and h >= 0.8:
                sizes = [
                    r.font.size.pt
                    for p in shape.text_frame.paragraphs
                    for r in p.runs
                    if r.font.size
                ]
                avg_pt = sum(sizes) / len(sizes) if sizes else 14
                # Rough glyph budget for the box at this size.
                capacity = (w * 72 / (avg_pt * 0.55)) * (h * 72 / (avg_pt * 1.5))
                if text_chars > capacity * 1.8:
                    flag(
                        "possible_overflow",
                        f"{text_chars} chars at ~{avg_pt:.0f}pt likely overflow "
                        f"a {w:.1f}x{h:.1f}in box",
                    )
    return findings


REVIEW_SCHEMA = {
    "type": "object",
    "properties": {
        "on_brand": {"type": "boolean"},
        "violations": {"type": "array", "items": {"type": "string"}},
    },
    "required": ["on_brand", "violations"],
    "additionalProperties": False,
}

REVIEW_PROMPT = """The FIRST {n_ref} image(s) are reference slides from the canonical Overview
brand deck. The LAST image is a generated slide under review.

Brand rules:
{rules}

Judge whether the reviewed slide belongs to the same visual family as the
references: palette (purple accent on clean light/dark grounds), typographic
hierarchy and weights, whitespace and margins, logo usage. Content differs
by design — judge only brand look. List concrete violations if any; minor
taste differences are not violations."""


def review_render(
    png_bytes: bytes,
    brand: dict | None = None,
    model: str | None = None,
    refs: list[bytes] | None = None,
) -> dict:
    """Vision brand review of one rendered slide. `refs` lets a caller read
    the reference renders once and reuse them across many reviews."""
    brand = brand or load_brand()
    refs = refs if refs is not None else [p.read_bytes() for p in reference_renders()[:3]]
    rules = yaml.safe_dump(
        {k: brand[k] for k in ("fonts", "colors", "typography", "notes") if k in brand},
        sort_keys=False,
    )
    try:
        return llm.complete(
            REVIEW_PROMPT.format(n_ref=len(refs), rules=rules),
            schema=REVIEW_SCHEMA,
            images=refs + [png_bytes],
            max_tokens=1000,
            model=model or llm.SONNET,
        )
    except llm.LLMRefusal:
        return {"on_brand": True, "violations": ["review refused; lint-only"]}


def _render_pages(pptx_path: Path) -> list[bytes]:
    """Every slide as PNG bytes, rasterized sequentially (local CPU work;
    a shared PyMuPDF document is not worth threading)."""
    from deck import render

    # FIDELITY: the audit judges how the deck looks, so it should judge the
    # rendering an engineer will actually see.
    pdf = render.pdf_bytes(pptx_path, purpose=render.FIDELITY)
    if pdf is None:
        return []
    import fitz

    doc = fitz.open(stream=pdf, filetype="pdf")
    try:
        return [doc[i].get_pixmap(dpi=72).tobytes("png") for i in range(len(doc))]
    finally:
        doc.close()


def generated_slide_numbers(included_slides: list[dict]) -> list[int]:
    """1-based positions (in the assembled deck) of slides this pipeline
    GENERATED — agent-built or freeform. Inherited skeleton slides are the
    company's own templates: linting them is useful, but judging their look
    against a handful of references produces false positives, so the vision
    tier is scoped to these."""
    # Presence, not truthiness: an entry can legitimately carry an empty
    # agent_slide spec and still be a generated slide.
    markers = ("agent_pptx", "agent_slide", "freeform")
    return [
        i
        for i, s in enumerate(included_slides, start=1)
        if any(k in s for k in markers)
    ]


def audit_deck(
    pptx_path: Path,
    report_path: Path,
    vision: bool = False,
    slides: list[int] | None = None,
    included_slides: list[dict] | None = None,
    log=print,
) -> dict:
    """Post-assembly, report-only brand audit -> brand_report.json.

    Lint covers every slide (objective and free). The vision tier is scoped:
    pass `slides` (1-based) explicitly, or `included_slides` (the plan's
    non-skipped entries) to have it default to the generated ones. Passing
    neither reviews the whole deck, which is rarely what you want.
    """
    if vision and slides is None and included_slides is not None:
        slides = generated_slide_numbers(included_slides)
        log(f"  brand vision audit scoped to {len(slides)} generated slide(s)")
    import shutil

    brand = load_brand()
    report: dict = {"deck": str(pptx_path), "lint": lint_presentation(pptx_path, brand)}
    if vision:
        reviews: list[dict] = []
        if not shutil.which("soffice"):
            reviews.append({"error": "LibreOffice not found; vision audit skipped"})
        else:
            try:
                pages = _render_pages(pptx_path)
            except Exception as e:
                pages = []
                reviews.append({"error": f"vision audit unavailable: {e}"})
            wanted = [
                (i, png)
                for i, png in enumerate(pages, start=1)
                if slides is None or i in set(slides)
            ]
            if wanted:
                # Reviews are independent; run them concurrently and merge in
                # slide order so the report stays deterministic.
                from concurrent.futures import ThreadPoolExecutor

                refs = [p.read_bytes() for p in reference_renders()[:3]]
                with ThreadPoolExecutor(max_workers=AUDIT_WORKERS) as pool:
                    futures = [
                        pool.submit(review_render, png, brand, None, refs)
                        for _i, png in wanted
                    ]
                    for (idx, _png), fut in zip(wanted, futures):
                        try:
                            verdict = fut.result()
                        except Exception as e:
                            reviews.append({"slide": idx, "error": f"review failed: {e}"})
                            continue
                        if not verdict["on_brand"]:
                            reviews.append({"slide": idx, **verdict})
        report["vision"] = reviews
    report_path.write_text(json.dumps(report, indent=2))
    n = len(report["lint"])
    v = len([r for r in report.get("vision", []) if "slide" in r])
    log(
        f"  brand audit: {n} lint finding(s)"
        + (f", {v} vision finding(s)" if vision else "")
        + f" -> {report_path.name}"
    )
    return report
