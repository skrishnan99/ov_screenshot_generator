"""Deck assembly from per-slide skeleton .pptx files.

Strategy: every skeleton is a standalone one-slide presentation. Each slide is
FILLED in its own file with ordinary python-pptx operations (token replacement
in text frames, image swaps), and only then transplanted into the output deck
via OPC part surgery — the surgery layer never edits content, it only copies
finished parts. The first slide's skeleton serves as the base presentation, so
the deck inherits its masters/theme; appended slides bring their own layouts
(remapped onto the base master) and media with them.
"""

from __future__ import annotations

import os
import re
from pathlib import Path

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.oxml.ns import qn
from pptx.util import Emu

TOKEN_RE = re.compile(r"\{\{\s*([\w.]+)\s*\}\}")

MSO_PICTURE = 13
MSO_GROUP = 6

# Images are embedded at ~this density for the space they actually occupy.
# A 3840px capture in a 5.5in slot renders around 800px on screen, so
# carrying the native file makes the deck many times larger for no visible
# gain. 200 DPI still looks right when projected or printed. Run assets on
# disk are untouched — only what goes INTO the deck is resized.
EMBED_DPI = 200
_resized_cache: dict[tuple[str, int, int], str] = {}

# Content pictures vs decorations (sidebar strips, logos): a content slot is
# reasonably wide AND reasonably large.
MIN_PIC_WIDTH_IN = 1.0
MIN_PIC_AREA_SQIN = 2.0


def _encode_smaller(im) -> str:
    """Write the image as whichever of PNG/JPEG comes out smaller.

    These decks mix two very different kinds of picture: UI screenshots
    (flat colour and text, where PNG wins and JPEG would ring around the
    text) and camera captures of metal parts (photographs, where PNG is
    several times larger than a visually identical JPEG). Rather than guess
    from the file extension, encode both and keep the smaller — self-tuning
    and never worse than the source format. Transparency forces PNG.
    """
    import tempfile

    def _write(suffix: str, save):
        fd, path = tempfile.mkstemp(prefix="sg-embed-", suffix=suffix)
        os.close(fd)
        save(path)
        return path, os.path.getsize(path)

    png_path, png_size = _write(".png", lambda p: im.save(p, optimize=True))
    # Presence of an alpha channel is not the question — screenshots carry a
    # fully opaque one. Only genuinely transparent pixels force PNG.
    transparent = "transparency" in im.info
    if not transparent and im.mode in ("RGBA", "LA"):
        transparent = im.getchannel("A").getextrema()[0] < 255
    if transparent:
        return png_path
    jpg_path, jpg_size = _write(
        ".jpg", lambda p: im.convert("RGB").save(p, quality=90, optimize=True)
    )
    if jpg_size < png_size:
        os.unlink(png_path)
        return jpg_path
    os.unlink(jpg_path)
    return png_path


def sized_for_slot(image_path, width_emu: int, height_emu: int) -> str:
    """A copy of the image no larger than its slot needs at EMBED_DPI.
    Returns the original path when it is already small enough (or on any
    failure — a bigger deck beats a broken one). Results are cached, so the
    same asset used in several slots is resized once per size."""
    import tempfile

    from PIL import Image
    from pptx.util import Emu

    src = str(image_path)
    target_w = max(1, int(Emu(width_emu).inches * EMBED_DPI))
    target_h = max(1, int(Emu(height_emu).inches * EMBED_DPI))
    key = (src, target_w, target_h)
    if key in _resized_cache:
        return _resized_cache[key]
    try:
        with Image.open(src) as im:
            if im.width <= target_w and im.height <= target_h:
                _resized_cache[key] = src
                return src
            scale = min(target_w / im.width, target_h / im.height)
            out = im.convert("RGB") if im.mode in ("CMYK", "P") else im.copy()
            out = out.resize(
                (max(1, int(im.width * scale)), max(1, int(im.height * scale))),
                Image.LANCZOS,
            )
            tmp = _encode_smaller(out)
    except Exception:
        _resized_cache[key] = src
        return src
    _resized_cache[key] = tmp
    return tmp


def iter_shapes(container):
    for shape in container.shapes:
        if shape.shape_type == MSO_GROUP:
            yield from _iter_group(shape)
        else:
            yield shape


def _iter_group(group):
    for shape in group.shapes:
        if shape.shape_type == MSO_GROUP:
            yield from _iter_group(shape)
        else:
            yield shape


def find_tokens(slide) -> list[str]:
    tokens = []
    for shape in iter_shapes(slide):
        if shape.has_text_frame:
            tokens += TOKEN_RE.findall(shape.text_frame.text)
    return tokens


def fill_tokens(slide, values: dict) -> list[str]:
    """Replace {{ token }} occurrences; formatting of each line's first run
    is preserved. Runs are consolidated per LINE — between <a:br> soft
    breaks — never across a whole paragraph: Google exports split a token
    across runs *within* a line, while multi-line info blocks (the recipe
    title slide's Site/Project/Date/name box) are one paragraph whose soft
    breaks a paragraph-wide join would silently eat, running every line
    together. A paragraph without breaks is one segment, so single-line
    behavior is unchanged."""
    from pptx.oxml.ns import qn

    filled = []
    for shape in iter_shapes(slide):
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            segments: list[list] = [[]]
            for child in para._p:
                if child.tag == qn("a:br"):
                    segments.append([])
                elif child.tag == qn("a:r"):
                    segments[-1].append(child)
            for seg in segments:
                ts = [t for r in seg if (t := r.find(qn("a:t"))) is not None]
                if not ts:
                    continue
                text = "".join(t.text or "" for t in ts)
                hits = TOKEN_RE.findall(text)
                if not hits:
                    continue
                ts[0].text = TOKEN_RE.sub(
                    lambda m: str(values.get(m.group(1), m.group(0))), text
                )
                for t in ts[1:]:
                    t.text = ""
                filled += [h for h in hits if h in values]
    return filled


def content_pictures(slide) -> list:
    pics = []
    for shape in iter_shapes(slide):
        if shape.shape_type != MSO_PICTURE:
            continue
        w, h = Emu(shape.width).inches, Emu(shape.height).inches
        if w >= MIN_PIC_WIDTH_IN and w * h >= MIN_PIC_AREA_SQIN:
            pics.append(shape)
    pics.sort(key=lambda s: s.width * s.height, reverse=True)
    return pics


PLACEHOLDER_TEXT_RE = re.compile(r"insert screenshot", re.IGNORECASE)


def image_slots(slide) -> list:
    """Fillable image slots: content PICTURE shapes (largest first), then
    'Insert screenshot here' auto-shape placeholders (left-to-right)."""
    slots = content_pictures(slide)
    placeholders = [
        s
        for s in iter_shapes(slide)
        if s.shape_type != MSO_PICTURE
        and s.has_text_frame
        and PLACEHOLDER_TEXT_RE.search(s.text_frame.text)
    ]
    placeholders.sort(key=lambda s: (s.left, s.top))
    return slots + placeholders


def _fill_placeholder(slide, shape, image_path: Path) -> None:
    """Replace an 'Insert screenshot here' auto-shape with the image,
    aspect-fit within the shape's bounds."""
    from PIL import Image

    with Image.open(image_path) as im:
        iw, ih = im.size
    scale = min(shape.width / iw, shape.height / ih)
    w, h = int(iw * scale), int(ih * scale)
    left = shape.left + (shape.width - w) // 2
    top = shape.top + (shape.height - h) // 2
    slide.shapes.add_picture(sized_for_slot(image_path, w, h), left, top, w, h)
    shape._element.getparent().remove(shape._element)


def fill_images(slide, image_paths: list) -> None:
    slots = image_slots(slide)
    if len(image_paths) > len(slots):
        raise RuntimeError(
            f"{len(image_paths)} images provided but only {len(slots)} slots on slide"
        )
    for path, slot in zip(image_paths, slots):
        if slot.shape_type == MSO_PICTURE:
            replace_picture(slide, slot, Path(path))
        else:
            _fill_placeholder(slide, slot, Path(path))
    # An unfilled placeholder must not ship its "Insert screenshot" note.
    for slot in slots[len(image_paths):]:
        if slot.shape_type != MSO_PICTURE and slot.has_text_frame:
            for para in slot.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""


def replace_picture(slide, pic, image_path: Path) -> None:
    """Swap a picture's image, aspect-fitting it within the original bounds."""
    from PIL import Image

    with Image.open(image_path) as im:
        iw, ih = im.size
    left, top, width, height = pic.left, pic.top, pic.width, pic.height
    scale = min(width / iw, height / ih)
    new_w, new_h = int(iw * scale), int(ih * scale)
    pic.left = left + (width - new_w) // 2
    pic.top = top + (height - new_h) // 2
    pic.width, pic.height = new_w, new_h

    image_part, rId = pic.part.get_or_add_image_part(
        sized_for_slot(image_path, new_w, new_h)
    )
    blip = pic._element.blipFill.find(qn("a:blip"))
    blip.set(qn("r:embed"), rId)


# ---------------------------------------------------------------------------
# Freeform slides: no dedicated skeleton. A donor skeleton supplies the
# theme/background/decorations; its token text and image slots are stripped
# and title/body/images are laid out programmatically.
# ---------------------------------------------------------------------------


def fill_freeform(pres: Presentation, job: dict) -> None:
    from pptx.util import Inches, Pt

    from deck.brand import load_brand

    brand = load_brand()
    slide = pres.slides[0]
    doomed = set()
    for shape in iter_shapes(slide):
        if shape.has_text_frame and TOKEN_RE.search(shape.text_frame.text):
            doomed.add(shape._element)
    for slot in image_slots(slide):
        doomed.add(slot._element)
    for el in doomed:
        el.getparent().remove(el)

    tokens = job.get("tokens") or {}
    title = tokens.get("_ff_title", "")
    body = tokens.get("_ff_body", "")
    images = job.get("images") or []
    sw, sh = pres.slide_width, pres.slide_height
    margin = Inches(0.55)

    # Geometry below is the numbered configuration step from
    # deck/brand/design_guide.md — "byte-for-byte identical" across the five
    # step templates, and the guide's instruction is to match it exactly.
    # This layout is the fallback for a failed agent slide, so it lands in
    # the middle of the numbered run and has to read as one of them; the
    # previous ad-hoc geometry (text left, image right, 22 pt title, grey
    # 13 pt body) was conspicuously not one of them.
    if title:
        from pptx.dml.color import RGBColor

        tb = slide.shapes.add_textbox(
            Inches(0.74), Inches(0.42), Inches(8.63), Inches(0.50)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        para = tf.paragraphs[0]
        para.line_spacing = 1.20
        run = para.add_run()
        run.text = title
        run.font.name = brand["fonts"]["headline"]
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(
            brand["colors"]["headline_black"].lstrip("#")
        )
    content_top = Inches(1.25)
    content_h = sh - content_top - Inches(0.35)
    if body and images:
        # Left two-thirds evidence, right third explanation.
        _freeform_images(
            slide, images, Inches(0.70), Inches(1.73), Inches(5.48), Inches(3.14)
        )
        # Height 0.34 GROWING DOWNWARD, per the guide — not a tall fixed box.
        # The fixed 2.85 top optically centres a short block against the image
        # band; a full-height box would start at 2.85 and run off the bottom
        # of the slide (brand lint: "text shape ... exceeds the slide").
        _freeform_body(
            slide, body, Inches(6.30), Inches(2.85), Inches(3.60), Inches(0.34),
            size_pt=17, bold=True, color=brand["colors"]["primary_purple"],
            grow=True,
        )
    elif images:
        _freeform_images(slide, images, margin, content_top, sw - 2 * margin, content_h)
    elif body:
        _freeform_body(slide, body, margin, content_top, sw - 2 * margin, content_h)


def _freeform_body(
    slide, text: str, left, top, width, height,
    size_pt: int = 13, bold: bool = False, color: str | None = None,
    grow: bool = False,
) -> None:
    """Body copy block. Defaults suit a text-only slide; the numbered-step
    right-hand column passes Bold 17 pt purple per the design guide.

    ``grow`` declares a short box that expands to fit its text, which is how
    the corpus authors the right-hand column. The stored height stays small,
    so the shape's declared extent remains inside the slide.
    """
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.util import Pt

    from deck.brand import load_brand

    brand = load_brand()
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    if grow:
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    for i, line in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.line_spacing = 1.0
        run = para.add_run()
        run.text = line
        run.font.name = brand["fonts"]["body"]
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(
            (color or brand["colors"]["body_dark"]).lstrip("#")
        )


def _freeform_images(slide, paths: list, left, top, width, height) -> None:
    """Stack images vertically, each aspect-fit within its cell."""
    from pptx.util import Inches
    from PIL import Image

    gap = Inches(0.15)
    cell_h = int((height - gap * (len(paths) - 1)) / len(paths))
    for i, path in enumerate(paths):
        with Image.open(path) as im:
            iw, ih = im.size
        cell_top = top + i * (cell_h + gap)
        scale = min(width / iw, cell_h / ih)
        w, h = int(iw * scale), int(ih * scale)
        slide.shapes.add_picture(
            sized_for_slot(path, w, h),
            left + (width - w) // 2,
            cell_top + (cell_h - h) // 2,
            w,
            h,
        )


# ---------------------------------------------------------------------------
# Theme baking: make a slide's colours independent of which theme wins.
# ---------------------------------------------------------------------------

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _theme_colors(slide) -> dict:
    """Scheme-colour name → literal hex, resolved from the slide's OWN theme.

    Includes the bg1/tx1/bg2/tx2 aliases, which are not scheme entries but
    indirections through the master's ``clrMap`` — and which a layout or
    slide may flip via ``clrMapOvr`` (dark designs swap the text and
    background roles). Missing that override renders light text as dark.
    Returns {} on any problem, which makes baking a no-op.
    """
    from lxml import etree
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    try:
        master = slide.slide_layout.slide_master
        theme = etree.fromstring(master.part.part_related_by(RT.THEME).blob)
        colors: dict[str, str] = {}
        scheme = theme.find(f".//{{{_A_NS}}}clrScheme")
        for child in scheme if scheme is not None else []:
            name = etree.QName(child).localname
            srgb = child.find(f"{{{_A_NS}}}srgbClr")
            sysc = child.find(f"{{{_A_NS}}}sysClr")
            val = (
                srgb.get("val") if srgb is not None
                else sysc.get("lastClr") if sysc is not None
                else None
            )
            if val:
                colors[name] = val

        aliases: dict[str, str] = {}
        clr_map = master.element.find(f"{{{_P_NS}}}clrMap")
        if clr_map is not None:
            aliases.update(clr_map.attrib)
        for owner in (slide.slide_layout, slide):
            ovr = owner._element.find(
                f"{{{_P_NS}}}clrMapOvr/{{{_A_NS}}}overrideClrMapping"
            )
            if ovr is not None:
                aliases.update(ovr.attrib)
        for alias, target in aliases.items():
            if target in colors:
                colors[alias] = colors[target]
        return colors
    except Exception:
        return {}


def bake_theme_colors(slide) -> int:
    """Rewrite ``<a:schemeClr val="accent1">`` as ``<a:srgbClr val="…">``
    using the slide's own theme, so the slide carries its colours rather
    than referencing a theme that a merge or an importer might change.

    Only colours: our skeletons use explicit typefaces (zero ``+mn-lt`` /
    ``+mj-lt`` tokens), which is why fonts already survive conversion.
    Transform children (lumMod, alpha, …) are preserved automatically —
    ``srgbClr`` accepts the same children as ``schemeClr``. Unknown
    references are left alone. Returns how many references were baked.
    """
    colors = _theme_colors(slide)
    if not colors:
        return 0
    baked = 0
    for sc in slide._element.iter(f"{{{_A_NS}}}schemeClr"):
        hex_val = colors.get(sc.get("val"))
        if hex_val:
            sc.tag = f"{{{_A_NS}}}srgbClr"
            sc.set("val", hex_val)
            baked += 1
    return baked


# ---------------------------------------------------------------------------
# OPC surgery: transplant a finished slide into the base presentation.
# ---------------------------------------------------------------------------


def _add_rel(part, reltype, target, rId, external=False):
    """Add a relationship preserving the ORIGINAL rId — the transplanted
    slide XML references relationships by rId (r:embed etc.), so auto-assigned
    ids would break it. Mirrors _Relationships.load_from_xml's construction."""
    from pptx.opc.oxml import RTM
    from pptx.opc.package import _Relationship

    rels = part.rels
    rels._rels[rId] = _Relationship(
        rels._base_uri,
        rId,
        reltype,
        target_mode=RTM.EXTERNAL if external else RTM.INTERNAL,
        target=target,
    )


def _partname_template(src_partname: str) -> str:
    if "/slideLayouts/" in src_partname:
        return "/ppt/slideLayouts/slideLayout%d.xml"
    if "/slides/" in src_partname:
        return "/ppt/slides/slide%d.xml"
    if "/media/" in src_partname:
        ext = src_partname.rsplit(".", 1)[-1]
        return f"/ppt/media/image%d.{ext}"
    if "/notesSlides/" in src_partname:
        return "/ppt/notesSlides/notesSlide%d.xml"
    stem, ext = src_partname.rsplit(".", 1)
    stem = re.sub(r"\d+$", "", stem)
    return f"{stem}%d.{ext}"


def _register_master(base: Presentation, master_part) -> None:
    """Make an imported slide master part official: relate it to the
    presentation part and add a sldMasterId entry."""
    pres_part = base.part
    rId = pres_part.relate_to(master_part, RT.SLIDE_MASTER)
    pres_el = pres_part._element
    lst = pres_el.find(qn("p:sldMasterIdLst"))
    existing = [int(e.get("id")) for e in lst.findall(qn("p:sldMasterId"))]
    el = lst.makeelement(qn("p:sldMasterId"), {"id": str(max(existing) + 1)})
    el.set(qn("r:id"), rId)
    lst.append(el)


def _alloc_partname(tmpl: str, used: set):
    from pptx.opc.packuri import PackURI

    n = 1
    while tmpl % n in used:
        n += 1
    used.add(tmpl % n)
    return PackURI(tmpl % n)


def _import_part(base: Presentation, src_part, mapping: dict, ctx: dict):
    key = str(src_part.partname)
    if key in mapping:
        return mapping[key]
    ct = src_part.content_type
    blob = src_part.blob
    # Identical parts (shared layouts, the decorative sidebar image every
    # skeleton carries) are stored once per deck.
    content_key = (ct, hash(blob))
    if content_key in ctx["by_hash"]:
        mapping[key] = ctx["by_hash"][content_key]
        return mapping[key]
    pkg = base.part.package
    partname = _alloc_partname(_partname_template(str(src_part.partname)), ctx["used"])
    new_part = Part(partname, ct, pkg, blob)
    mapping[key] = new_part
    ctx["by_hash"][content_key] = new_part
    for rId, rel in sorted(src_part.rels.items()):
        if rel.is_external:
            _add_rel(new_part, rel.reltype, rel.target_ref, rId, external=True)
            continue
        if rel.target_part.content_type.endswith("notesSlide+xml"):
            continue  # notes are irrelevant to the deck; slide XML holds no rId to them
        target = _import_part(base, rel.target_part, mapping, ctx)
        _add_rel(new_part, rel.reltype, target, rId)
    if ct.endswith("slideMaster+xml"):
        # Slides keep their own master (backgrounds/themes differ per
        # skeleton); identical masters dedup via by_hash above, so the shared
        # content master is imported and registered exactly once.
        _register_master(base, new_part)
    return new_part


def append_slide(
    base: Presentation, src_pres: Presentation, ctx: dict, index: int = 0
) -> None:
    """Append one of src's (already-filled) slides to the base deck.

    Defaults to the first slide, which is what the skeleton pipeline builds.
    ``index`` lets a caller lift a specific slide out of a multi-slide source —
    used to transplant boilerplate verbatim from the reference template rather
    than re-authoring an approximation of it.
    """
    src_slide_part = src_pres.slides[index].part
    mapping: dict = {}
    new_part = _import_part(base, src_slide_part, mapping, ctx)
    rId = base.part.relate_to(new_part, RT.SLIDE)
    base.slides._sldIdLst.add_sldId(rId)


def build_deck(slide_jobs: list[dict], out_path: Path) -> None:
    """slide_jobs: [{skeleton: Path, tokens: {name: value}, image: Path|None}].
    The first job's skeleton becomes the base presentation."""
    if not slide_jobs:
        raise ValueError("no slides to build")

    def load(job):
        # Agent-built slides arrive as finished one-slide files; everything
        # else starts from its skeleton and gets filled below.
        if job.get("agent_pptx"):
            return Presentation(str(job["agent_pptx"]))
        return Presentation(str(job["skeleton"]))

    def fill(pres, job):
        if job.get("agent_pptx"):
            return  # authored complete by the agent session
        if job.get("freeform"):
            fill_freeform(pres, job)
            return
        slide = pres.slides[0]
        if job.get("tokens"):
            fill_tokens(slide, job["tokens"])
        images = job.get("images") or ([job["image"]] if job.get("image") else [])
        if images:
            fill_images(slide, images)

    base = load(slide_jobs[0])
    fill(base, slide_jobs[0])
    # Bake AFTER filling (filled runs can carry scheme colours too) and
    # before transplant, while the slide still sits with its own theme.
    bake_theme_colors(base.slides[0])
    ctx = {
        "used": {str(p.partname) for p in base.part.package.iter_parts()},
        "by_hash": {},
    }
    for job in slide_jobs[1:]:
        src = load(job)
        fill(src, job)
        bake_theme_colors(src.slides[0])
        append_slide(base, src, ctx)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    base.save(str(out_path))
