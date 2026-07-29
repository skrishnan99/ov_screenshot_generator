"""Skeleton introspection + sidecar content-hole descriptions.

A skeleton's content holes are described in an optional sidecar YAML sitting
next to the .pptx with the same stem (imaging_setup.pptx -> imaging_setup.yaml).
The sidecar is authored metadata — it never renders — and says what each hole
should CONTAIN, not which file fills it:

    image_slots:            # aligned with the skeleton's fill order
      - name: main_screen   # optional, for readable reports
        expects: >
          The Imaging Setup screen showing ...
    tokens:                 # optional default guidance for text tokens the
      setup_text: >         # deck spec does not cover itself
        1-2 sentences about ...

The profile merges what the pptx actually contains (discovered tokens and
image slots, slide title, slot geometry) with the sidecar, and reports drift
between the two as warnings instead of failing silently.
"""

from __future__ import annotations

from functools import lru_cache
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.util import Emu

from deck.assemble import MSO_PICTURE, find_tokens, image_slots, iter_shapes


def _slide_title(slide) -> str:
    """Topmost non-trivial text on the slide — the headline, used as context
    when describing what the slide (and thus its holes) is about."""
    best = None
    for shape in iter_shapes(slide):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text or text == "‹#›":
            continue
        if best is None or shape.top < best[0]:
            best = (shape.top, text)
    return best[1].replace("\n", " / ") if best else ""


def _load_sidecar(pptx_path: Path) -> dict:
    sidecar_path = pptx_path.with_suffix(".yaml")
    if not sidecar_path.exists():
        return {}
    data = yaml.safe_load(sidecar_path.read_text()) or {}
    if not isinstance(data, dict):
        raise ValueError(f"sidecar {sidecar_path} must be a YAML mapping")
    return data


@lru_cache(maxsize=None)
def skeleton_profile(pptx_path: str) -> dict:
    """Everything the matcher needs to know about a skeleton's content holes:

    {
      "title": slide headline,
      "tokens": [token names present in the pptx],
      "token_guidance": {token: sidecar guidance},
      "slots": [{index, name, is_picture, width_in, height_in, aspect,
                 placeholder_text, expects}],   # in fill order
      "warnings": [drift messages],
    }
    """
    path = Path(pptx_path)
    pres = Presentation(str(path))
    slide = pres.slides[0]
    sidecar = _load_sidecar(path)
    warnings: list[str] = []

    tokens = find_tokens(slide)
    slots = []
    for i, shape in enumerate(image_slots(slide)):
        placeholder_text = ""
        if shape.shape_type != MSO_PICTURE and shape.has_text_frame:
            placeholder_text = " ".join(shape.text_frame.text.split())
        w, h = Emu(shape.width).inches, Emu(shape.height).inches
        slots.append(
            {
                "index": i,
                "name": f"slot_{i}",
                "is_picture": shape.shape_type == MSO_PICTURE,
                "width_in": round(w, 2),
                "height_in": round(h, 2),
                "aspect": round(w / h, 2) if h else 0,
                "placeholder_text": placeholder_text,
                "expects": "",
            }
        )

    side_slots = sidecar.get("image_slots") or []
    if side_slots and len(side_slots) != len(slots):
        warnings.append(
            f"{path.name}: sidecar describes {len(side_slots)} image slot(s) "
            f"but the skeleton has {len(slots)} — check for drift"
        )
    for slot, side in zip(slots, side_slots):
        if side.get("name"):
            slot["name"] = side["name"]
        if side.get("expects"):
            slot["expects"] = " ".join(str(side["expects"]).split())

    token_guidance = {}
    for name, guidance in (sidecar.get("tokens") or {}).items():
        if name not in tokens:
            warnings.append(
                f"{path.name}: sidecar describes token '{name}' that does not "
                f"exist in the skeleton — check for drift"
            )
            continue
        token_guidance[name] = " ".join(str(guidance).split())

    return {
        "path": str(path),
        "title": _slide_title(slide),
        "tokens": tokens,
        "token_guidance": token_guidance,
        "slots": slots,
        "warnings": warnings,
    }
