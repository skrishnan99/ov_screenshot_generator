#!/usr/bin/env python3
"""ovdeck — the Overview.ai branded deck engine.

One import gives you the brand palette, a fixed set of validated layouts, and a
self-checking build. You describe slides; the engine does every placement
calculation, measures every string against real font metrics, and refuses to
save a deck that has text overflowing its box, shapes off the canvas, colliding
content, or a colour that is not in the brand palette.

    from ovdeck import Deck

    d = Deck("out/report.pptx")                    # style="report"
    d.title_slide("OV80i AI Vision Inspection", "Acme Ltd - Line 3 Demo",
                  meta=["Report by: Jordan Lee", "Date: 2026.07.30"],
                  image="assets/part.jpg")
    d.contents([("01", "Introduction", "What was demonstrated")])
    d.section("01", "Introduction")
    d.figure("Step 1: Imaging Setup", "screenshots/imaging.png",
             caption="Exposure and lighting tuned live on the part.",
             chips=["3840x2160", "Exposure 20 ms"])
    d.closing(para="Thank you.", summary=["OV80i - serial gsac177082"])
    d.save()          # raises LayoutError if anything is wrong

Two house styles, both taken from the example decks:

    style="report"        navy header bar, tinted page   (STADLER connector deck)
    style="presentation"  white page, purple left spine  (Hot Bar Soldering deck)

Layouts are identical between them; only the chrome changes. Every layout is
documented in references/layouts.md. Never position shapes by hand outside this
module — that is the whole point of it.
"""
from __future__ import annotations

import glob
import json
import os
from dataclasses import dataclass, field
from pathlib import Path
from typing import Sequence

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
SKILL = HERE.parent
TOKENS = json.loads((SKILL / "assets" / "tokens.json").read_text())
BRAND_DIR = SKILL / "assets" / "brand"
DERIVED = BRAND_DIR / "derived"

# --------------------------------------------------------------------------
# brand tokens
# --------------------------------------------------------------------------


def _rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str.lstrip("#").upper())


_SEM = TOKENS["semantic"]

PAGE_BG = _rgb(_SEM["page_bg"])
PAGE_BG_TINT = _rgb(_SEM["page_bg_tinted"])
SURFACE = _rgb(_SEM["surface"])
SURFACE_ALT = _rgb(_SEM["surface_alt"])
HEADER_BAR = _rgb(_SEM["header_bar"])
DARK_BG = _rgb(_SEM["dark_bg"])
DARK_BG_ALT = _rgb(_SEM["dark_bg_alt"])
ACCENT = _rgb(_SEM["accent"])
ACCENT_LIGHT = _rgb(_SEM["accent_light"])
ACCENT_ON_DARK = _rgb(_SEM["accent_on_dark"])
ACCENT_SOFT = _rgb(_SEM["accent_soft"])
CHIP_BG = _rgb(_SEM["chip_bg"])
CHIP_BORDER = _rgb(_SEM["chip_border"])
HAIRLINE = _rgb(_SEM["hairline"])
TEXT_PRIMARY = _rgb(_SEM["text_primary"])
TEXT_BODY = _rgb(_SEM["text_body"])
TEXT_MUTED = _rgb(_SEM["text_muted"])
TEXT_ON_DARK = _rgb(_SEM["text_on_dark"])
TEXT_ON_DARK_MUTED = _rgb(_SEM["text_on_dark_muted"])
HIGHLIGHT = _rgb(_SEM["highlight"])
HIGHLIGHT_SOFT = _rgb(_SEM["highlight_soft"])
WHITE = _rgb("#FFFFFF")

ALLOWED_COLORS = (
    {str(_rgb(h)) for ramp in TOKENS["ramps"].values() for h in ramp}
    | {str(_rgb(TOKENS["core"][k]["hex"])) for k in TOKENS["core"]}
    | {str(_rgb(v)) for v in TOKENS["neutral"].values()}
    | {str(_rgb(v)) for v in _SEM.values()}
)

DEFAULT_FONT = TOKENS["typography"]["family"]
PREFERRED_FONT = TOKENS["typography"].get("preferred_family")
SZ = TOKENS["typography"]["scale_pt"]

SLIDE_W, SLIDE_H = TOKENS["geometry"]["slide_in"]
MARGIN = TOKENS["geometry"]["margin_in"]
HEADER_H = TOKENS["geometry"]["header_bar_h_in"]
CONTENT_W = SLIDE_W - 2 * MARGIN
GUTTER = TOKENS["geometry"]["gutter_in"]
BODY_BOTTOM = 7.05          # nothing but footnotes may go below this
FOOTNOTE_Y = 6.84
SPINE_W = 0.55              # presentation style only


class LayoutError(RuntimeError):
    """Raised at save() when the deck violates a layout or brand rule."""


@dataclass
class Issue:
    slide: int
    kind: str
    detail: str
    severity: str = "error"

    def __str__(self) -> str:
        return f"[slide {self.slide:>2}] {self.severity.upper():7s} {self.kind}: {self.detail}"


@dataclass
class Rect:
    x: float
    y: float
    w: float
    h: float
    kind: str = "shape"
    label: str = ""

    @property
    def right(self) -> float:
        return self.x + self.w

    @property
    def bottom(self) -> float:
        return self.y + self.h

    def overlap_area(self, other: "Rect") -> float:
        dx = min(self.right, other.right) - max(self.x, other.x)
        dy = min(self.bottom, other.bottom) - max(self.y, other.y)
        return dx * dy if dx > 0 and dy > 0 else 0.0


# --------------------------------------------------------------------------
# text metrics — real font measurement, not guesswork
# --------------------------------------------------------------------------

_METRIC_SEARCH = [
    "/Applications/LibreOffice.app/Contents/Resources/fonts/truetype/Carlito-{w}.ttf",
    "/usr/share/fonts/truetype/crosextra/Carlito-{w}.ttf",
    "/Library/Fonts/Microsoft/Calibri{m}.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans{d}.ttf",
]
_FONT_DIRS = [
    "/Library/Fonts", "/System/Library/Fonts", "/System/Library/Fonts/Supplemental",
    os.path.expanduser("~/Library/Fonts"), "/usr/share/fonts", "/usr/local/share/fonts",
]
_font_cache: dict[tuple[float, bool], object] = {}
PX_PER_IN = 96.0
PT_TO_PX = PX_PER_IN / 72.0


def font_installed(family: str) -> bool:
    """Is this typeface actually available on this machine?"""
    key = family.replace(" ", "").lower()
    for d in _FONT_DIRS:
        for ext in ("ttf", "otf", "ttc"):
            for p in glob.glob(os.path.join(d, f"**/*.{ext}"), recursive=True):
                if os.path.basename(p).replace(" ", "").lower().startswith(key):
                    return True
    return False


def _metric_file(bold: bool) -> str | None:
    for tpl in _METRIC_SEARCH:
        p = tpl.format(w="Bold" if bold else "Regular",
                       m="-Bold" if bold else "",
                       d="-Bold" if bold else "")
        if os.path.exists(p):
            return p
    return None


def _pil_font(size_pt: float, bold: bool):
    key = (round(size_pt, 2), bold)
    if key in _font_cache:
        return _font_cache[key]
    from PIL import ImageFont

    path = _metric_file(bold)
    try:
        f = ImageFont.truetype(path, int(round(size_pt * PT_TO_PX))) if path else ImageFont.load_default()
    except Exception:                                          # pragma: no cover
        f = ImageFont.load_default()
    _font_cache[key] = f
    return f


def text_width_in(s: str, size_pt: float, bold: bool = False) -> float:
    """Rendered width of a single line, in inches."""
    f = _pil_font(size_pt, bold)
    try:
        return float(f.getlength(s)) / PX_PER_IN
    except AttributeError:                                     # pragma: no cover
        return len(s) * size_pt * 0.5 / 72.0


def wrap_lines(s: str, width_in: float, size_pt: float, bold: bool = False) -> list[str]:
    """Greedy word wrap using real glyph advances."""
    if not s:
        return [""]
    words, lines, cur = s.split(), [], ""
    for word in words:
        trial = f"{cur} {word}".strip()
        if text_width_in(trial, size_pt, bold) <= width_in or not cur:
            cur = trial
        else:
            lines.append(cur)
            cur = word
    if cur:
        lines.append(cur)
    return lines


def text_height_in(s: str, width_in: float, size_pt: float, bold: bool = False,
                   line_spacing: float = 1.0) -> float:
    """Height a wrapped paragraph will occupy, in inches (PowerPoint metrics)."""
    n = len(wrap_lines(s, width_in, size_pt, bold))
    return n * (size_pt * 1.2 * line_spacing) / 72.0


# --------------------------------------------------------------------------
# deck
# --------------------------------------------------------------------------


@dataclass
class _SlideState:
    index: int
    dark: bool
    rects: list[Rect] = field(default_factory=list)
    has_logo: bool = False
    title: str = ""


class Deck:
    """A brand-locked Overview.ai deck.

    style="report"        navy header bar on a tinted page (default)
    style="presentation"  white page with a purple left spine
    font=None             use the token default; pass a family to request it —
                          if it is not installed you get a warning and a
                          metric-safe fallback rather than silent substitution
    """

    def __init__(self, out_path: str | Path, *, strict: bool = True,
                 style: str = "report", font: str | None = None):
        if style not in ("report", "presentation"):
            raise ValueError("style must be 'report' or 'presentation'")
        self.out_path = Path(out_path)
        self.strict = strict
        self.style = style

        self.font = DEFAULT_FONT
        self._font_note: str | None = None
        if font and font != DEFAULT_FONT:
            if font_installed(font):
                self.font = font
            else:
                self._font_note = (
                    f"'{font}' is not installed on this machine; using {DEFAULT_FONT}. "
                    "Install the font or build on a machine that has it if the "
                    "customer copy must use it."
                )

        if style == "report":
            self.page_bg = PAGE_BG_TINT
            self.ML = MARGIN
            self.title_h = HEADER_H
        else:
            self.page_bg = PAGE_BG
            self.ML = SPINE_W + 0.50
            self.title_h = 1.42
        self.CW = SLIDE_W - self.ML - MARGIN

        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_W)
        self.prs.slide_height = Inches(SLIDE_H)
        self._blank = self.prs.slide_layouts[6]
        self.states: list[_SlideState] = []
        self.issues: list[Issue] = []
        self._colors_used: set[str] = set()
        self._fonts_used: set[str] = set()

    # ---------------- primitives ----------------

    def _new(self, dark: bool = False) -> tuple:
        sl = self.prs.slides.add_slide(self._blank)
        st = _SlideState(index=len(self.states) + 1, dark=dark)
        self.states.append(st)
        self._ends_with_template = False
        self._rect(sl, st, 0, 0, SLIDE_W, SLIDE_H,
                   fill=DARK_BG if dark else self.page_bg, register=False)
        if not dark and self.style == "presentation":
            self._rect(sl, st, 0, 0, SPINE_W, SLIDE_H, fill=ACCENT, register=False)
            self._text(sl, st, 0, 0.42, SPINE_W, 0.3,
                       [{"text": str(st.index), "size": SZ["footnote"], "bold": True,
                         "color": TEXT_ON_DARK, "spacing": 1.0}],
                       align=PP_ALIGN.CENTER, check=False)
            mark = DERIVED / "logomark.png"
            if mark.exists():
                iw, ih = Image.open(mark).size
                w = 0.30
                sl.shapes.add_picture(str(mark), Inches((SPINE_W - w) / 2),
                                      Inches(SLIDE_H - 0.62), Inches(w),
                                      Inches(w * ih / iw))
        return sl, st

    def _note(self, st: _SlideState | None, kind: str, detail: str,
              severity: str = "error") -> None:
        self.issues.append(Issue(st.index if st else 0, kind, detail, severity))

    def _reg(self, st: _SlideState, r: Rect) -> None:
        st.rects.append(r)
        if r.x < -0.01 or r.y < -0.01 or r.right > SLIDE_W + 0.01 or r.bottom > SLIDE_H + 0.01:
            self._note(st, "off-canvas",
                       f"{r.kind} '{r.label}' at ({r.x:.2f},{r.y:.2f}) "
                       f"{r.w:.2f}x{r.h:.2f} extends past the slide")

    def _rect(self, sl, st, x, y, w, h, *, fill=None, line=None, line_w=0.75,
              shape=MSO_SHAPE.RECTANGLE, register=True, label="") -> object:
        s = sl.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        if fill is None:
            s.fill.background()
        else:
            s.fill.solid()
            s.fill.fore_color.rgb = fill
            self._colors_used.add(str(fill))
        if line is None:
            s.line.fill.background()
        else:
            s.line.color.rgb = line
            s.line.width = Pt(line_w)
            self._colors_used.add(str(line))
        s.shadow.inherit = False
        if register:
            self._reg(st, Rect(x, y, w, h, "surface", label))
        return s

    def _text(self, sl, st, x, y, w, h, paragraphs, *, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP, label="", check=True) -> float:
        """paragraphs: [{text,size,bold,color,spacing,space_after,space_before}]

        Returns the y coordinate just below the measured text.
        """
        tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        used = 0.0
        for i, p in enumerate(paragraphs):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.alignment = align
            spacing = p.get("spacing", 1.0)
            para.line_spacing = spacing
            sa, sb = p.get("space_after", 0), p.get("space_before", 0)
            if sa:
                para.space_after = Pt(sa)
            if sb:
                para.space_before = Pt(sb)
            run = para.add_run()
            run.text = p["text"]
            run.font.size = Pt(p["size"])
            run.font.bold = p.get("bold", False)
            run.font.color.rgb = p["color"]
            run.font.name = self.font
            self._colors_used.add(str(p["color"]))
            self._fonts_used.add(self.font)
            used += text_height_in(p["text"], w, p["size"], p.get("bold", False), spacing)
            used += (sa + sb) / 72.0
        if check:
            if used > h + 0.03:
                self._note(st, "text-overflow",
                           f"'{label or paragraphs[0]['text'][:38]}' needs {used:.2f}in "
                           f"but its box is {h:.2f}in")
            self._reg(st, Rect(x, y, w, min(max(used, 0.08), h), "text",
                               label or paragraphs[0]["text"][:32]))
        return y + used

    def _para(self, sl, st, x, y, w, s, *, size=None, color=TEXT_BODY, bold=False,
              spacing=1.22, max_h=2.2, align=PP_ALIGN.LEFT, label="") -> float:
        size = size or SZ["body"]
        h = min(max(text_height_in(s, w, size, bold, spacing) + 0.02, 0.2), max_h)
        return self._text(sl, st, x, y, w, h,
                          [{"text": s, "size": size, "bold": bold, "color": color,
                            "spacing": spacing}], align=align, label=label or s[:32])

    def _bullets(self, sl, st, x, y, w, items, *, size=None, color=TEXT_BODY,
                 gap=7, spacing=1.10, label="bullets") -> float:
        size = size or SZ["bullet"]
        paras = [{"text": f"•   {it}", "size": size, "bold": False, "color": color,
                  "spacing": spacing, "space_after": gap} for it in items]
        h = sum(text_height_in(p["text"], w, size, False, spacing) for p in paras)
        h += gap * len(paras) / 72.0 + 0.04
        return self._text(sl, st, x, y, w, h, paras, label=label)

    def _chips(self, sl, st, x, y, items, *, size=None, on_dark=False,
               width=None) -> float:
        """Row of spec chips, wrapping to further rows. Returns bottom y."""
        size = size or SZ["chip"]
        ch = 0.30
        cx, cy = x, y
        limit = x + (width if width is not None else self.CW)
        for it in items:
            w = text_width_in(it, size) + 0.34
            if cx + w > limit + 0.001 and cx > x:
                cx, cy = x, cy + ch + 0.12
            self._rect(sl, st, cx, cy, w, ch,
                       fill=DARK_BG_ALT if on_dark else CHIP_BG,
                       line=None if on_dark else CHIP_BORDER, label=f"chip:{it}")
            self._text(sl, st, cx, cy + 0.055, w, 0.22,
                       [{"text": it, "size": size, "bold": False,
                         "color": TEXT_ON_DARK if on_dark else TEXT_PRIMARY,
                         "spacing": 1.0}], align=PP_ALIGN.CENTER, check=False)
            cx += w + 0.12
        return cy + ch

    def _picture(self, sl, st, path, bx, by, bw, bh, *, border=True, label="") -> Rect:
        path = str(path)
        if not os.path.exists(path):
            self._note(st, "missing-image", f"{path} does not exist")
            return Rect(bx, by, bw, bh, "image", label)
        iw, ih = Image.open(path).size
        scale = min(bw / iw, bh / ih)
        w, h = iw * scale, ih * scale
        x, y = bx + (bw - w) / 2, by + (bh - h) / 2
        sl.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
        if border:
            self._rect(sl, st, x - 0.02, y - 0.02, w + 0.04, h + 0.04,
                       fill=None, line=HAIRLINE, register=False)
        r = Rect(x, y, w, h, "image", label or os.path.basename(path))
        self._reg(st, r)
        return r

    def _logo(self, sl, st, x, y, width, *, on_dark=True, vertical=False) -> None:
        name = ("logo-V-white.png" if vertical else "logo-H-white.png") if on_dark else \
               ("logo-V-dark.png" if vertical else "logo-H-dark.png")
        p = DERIVED / name
        if not p.exists():
            self._note(st, "missing-logo",
                       f"{p} missing — run scripts/make_logo_variants.py")
            return
        min_w = TOKENS["geometry"]["logo_min_width_in"]
        if width < min_w:
            self._note(st, "logo-too-small",
                       f"logo width {width:.2f}in is below the {min_w}in minimum")
        iw, ih = Image.open(p).size
        h = width * ih / iw
        sl.shapes.add_picture(str(p), Inches(x), Inches(y), Inches(width), Inches(h))
        self._reg(st, Rect(x, y, width, h, "logo", name))
        st.has_logo = True

    def _header(self, sl, st, title: str, subtitle: str = "") -> None:
        """Slide chrome + title. Returns nothing; content starts below title_h."""
        st.title = title
        size = SZ["slide_title"]
        avail = self.CW - 0.4
        while text_width_in(title, size, True) > avail and size > 15:
            size -= 0.5
        if self.style == "report":
            self._rect(sl, st, 0, 0, SLIDE_W, HEADER_H, fill=HEADER_BAR, register=False)
            self._text(sl, st, self.ML, 0.30, self.CW, 0.62,
                       [{"text": title, "size": size, "bold": True,
                         "color": TEXT_ON_DARK, "spacing": 1.0}],
                       label=f"title:{title[:28]}")
            if subtitle:
                self._note(st, "unsupported",
                           "report style has no room for a slide subtitle; "
                           "fold it into the caption", severity="warning")
        else:
            y = self._para(sl, st, self.ML, 0.34, self.CW, title, size=size, bold=True,
                           color=TEXT_PRIMARY, spacing=1.05, max_h=0.72,
                           label=f"title:{title[:28]}")
            if subtitle:
                self._para(sl, st, self.ML, y + 0.06, self.CW, subtitle,
                           size=SZ["lead"], bold=True, color=ACCENT, spacing=1.05,
                           max_h=0.42, label="subtitle")

    def _top(self) -> float:
        return self.title_h + 0.27

    # ---------------- layouts ----------------

    def title_slide(self, title: str, subtitle: str = "", *, meta: Sequence[str] = (),
                    image: str | None = None, footer: str = "www.overview.ai"):
        """Opening slide: navy field, white logo, headline, optional circular image."""
        sl, st = self._new(dark=True)
        if image:
            self._rect(sl, st, 8.05, 0.75, 4.70, 4.70, fill=ACCENT_LIGHT,
                       shape=MSO_SHAPE.OVAL, register=False)
            disc = _circle_crop(image, self.out_path.parent / "_title_disc.png")
            sl.shapes.add_picture(str(disc), Inches(8.45), Inches(1.15),
                                  Inches(3.90), Inches(3.90))
            self._reg(st, Rect(8.05, 0.75, 4.70, 4.70, "image", "title-disc"))
        self._logo(sl, st, 0.85, 0.70, 2.55, on_dark=True)
        right = 7.60 if image else SLIDE_W - MARGIN
        y = self._para(sl, st, 0.85, 2.15, right - 0.85, title, size=SZ["display"],
                       bold=True, color=TEXT_ON_DARK, spacing=1.12, max_h=2.3,
                       label="title-headline")
        if subtitle:
            y = self._para(sl, st, 0.85, y + 0.28, right - 0.85, subtitle,
                           size=SZ["lead"] + 1, color=ACCENT_ON_DARK, spacing=1.2,
                           label="title-subtitle")
            self._rect(sl, st, 0.85, y + 0.30, 1.75, 0.05, fill=ACCENT_LIGHT,
                       register=False)
            y += 0.35
        my = max(y + 0.60, 4.62)
        for i, line in enumerate(meta):
            self._para(sl, st, 0.85, my + i * 0.42, right - 0.85, line,
                       size=SZ["body"], color=TEXT_ON_DARK if i == 0 else ACCENT_ON_DARK,
                       spacing=1.0, label=f"meta{i}")
        if footer:
            self._para(sl, st, 0.85, 6.86, 4.0, footer, size=SZ["footnote"] - 0.5,
                       color=TEXT_ON_DARK, spacing=1.0, label="footer")

    def contents(self, items: Sequence[tuple[str, str, str]], *, heading="C O N T E N T S"):
        """Numbered table of contents on the navy field."""
        sl, st = self._new(dark=True)
        self._para(sl, st, 0.95, 0.78, 10.0, heading, size=SZ["display"] - 4,
                   bold=True, color=TEXT_ON_DARK, spacing=1.0, label="contents")
        n = len(items)
        top, step = (2.45, 1.42) if n <= 3 else (2.10, 1.18)
        for i, (num, t, sub) in enumerate(items):
            y = top + i * step
            self._para(sl, st, 0.95, y - 0.12, 1.5, num, size=SZ["slide_title"] + 5,
                       bold=True, color=ACCENT_LIGHT, spacing=1.0, label=f"num{num}")
            self._para(sl, st, 3.05, y - 0.05, 9.2, t, size=SZ["card_title"] + 2,
                       bold=True, color=TEXT_ON_DARK, spacing=1.0, label=f"toc{i}")
            if sub:
                self._para(sl, st, 3.05, y + 0.45, 9.2, sub, size=SZ["body"],
                           color=ACCENT_ON_DARK, spacing=1.0, label=f"tocsub{i}")

    def section(self, number: str, title: str, subtitle: str = ""):
        """Full-bleed navy section divider."""
        sl, st = self._new(dark=True)
        self._text(sl, st, 0, 2.10, SLIDE_W, 1.30,
                   [{"text": number, "size": SZ["section_number"], "bold": True,
                     "color": ACCENT_LIGHT, "spacing": 1.0}],
                   align=PP_ALIGN.CENTER, label=f"sec{number}")
        self._para(sl, st, 0, 3.78, SLIDE_W, title, size=SZ["slide_title"] + 2,
                   bold=True, color=TEXT_ON_DARK, spacing=1.1, align=PP_ALIGN.CENTER,
                   label="sec-title")
        if subtitle:
            self._para(sl, st, SLIDE_W / 2 - 4.5, 4.62, 9.0, subtitle, size=SZ["body"],
                       color=ACCENT_ON_DARK, spacing=1.15, align=PP_ALIGN.CENTER,
                       label="sec-sub")

    def cards(self, title: str, cards: Sequence[tuple[str, str]], *, columns: int = 3,
              subtitle: str = ""):
        """Benefit/feature grid. 2-6 cards; 3 columns by default."""
        sl, st = self._new()
        self._header(sl, st, title, subtitle)
        n = len(cards)
        if not 2 <= n <= 6:
            self._note(st, "card-count", f"{n} cards given; use 2-6")
        rows = (n + columns - 1) // columns
        cw = (self.CW - GUTTER * (columns - 1)) / columns
        top = self._top() + 0.14
        ch = (BODY_BOTTOM - top - GUTTER * (rows - 1)) / rows
        for i, (t, d) in enumerate(cards):
            cx = self.ML + (i % columns) * (cw + GUTTER)
            cy = top + (i // columns) * (ch + GUTTER)
            self._rect(sl, st, cx, cy, cw, ch, fill=SURFACE if self.style == "report"
                       else SURFACE_ALT, label=f"card{i}")
            self._rect(sl, st, cx + 0.28, cy + 0.30, 0.42, 0.42, fill=ACCENT_SOFT,
                       shape=MSO_SHAPE.OVAL, register=False)
            self._para(sl, st, cx + 0.88, cy + 0.32, cw - 1.16, t,
                       size=SZ["caption"] + 2, bold=True, color=TEXT_PRIMARY,
                       spacing=1.06, max_h=0.8, label=f"cardt{i}")
            self._para(sl, st, cx + 0.28, cy + 1.22, cw - 0.56, d,
                       size=SZ["caption"], color=TEXT_BODY, spacing=1.16,
                       max_h=ch - 1.35, label=f"cardd{i}")

    def statement(self, title: str, intro: str, *, card_title: str,
                  bullets: Sequence[str], badge: str = "", rule: bool = True):
        """Intro paragraph plus one emphasised card with a bullet list."""
        sl, st = self._new()
        self._header(sl, st, title)
        y = self._para(sl, st, self.ML, self._top(), self.CW, intro, size=SZ["body"],
                       color=TEXT_BODY, spacing=1.25, max_h=1.1, label="intro")
        cy = y + 0.36
        cx = self.ML + 0.72
        cw = self.CW - 1.44
        ch = BODY_BOTTOM - cy
        if rule:
            self._rect(sl, st, cx, cy, cw, 0.09, fill=ACCENT, register=False)
        self._rect(sl, st, cx, cy + 0.09, cw, ch - 0.09,
                   fill=SURFACE if self.style == "report" else SURFACE_ALT,
                   label="stmt-card")
        ty = cy + 0.50
        tw = cw - 1.2 - (2.2 if badge else 0.0)
        self._para(sl, st, cx + 0.55, ty, tw, card_title, size=SZ["card_title"] + 1,
                   bold=True, color=TEXT_PRIMARY, spacing=1.06, max_h=0.6,
                   label="card-title")
        if badge:
            bw = text_width_in(badge, SZ["caption"], True) + 0.5
            self._rect(sl, st, cx + cw - 0.55 - bw, ty - 0.04, bw, 0.44,
                       fill=HIGHLIGHT_SOFT, line=HIGHLIGHT, label="badge")
            self._text(sl, st, cx + cw - 0.55 - bw, ty + 0.07, bw, 0.28,
                       [{"text": badge, "size": SZ["caption"], "bold": True,
                         "color": TEXT_PRIMARY, "spacing": 1.0}],
                       align=PP_ALIGN.CENTER, check=False)
        self._bullets(sl, st, cx + 0.60, ty + 0.80, cw - 1.2, bullets,
                      size=SZ["body"], gap=9)

    def figure(self, title: str, image: str, *, caption: str = "",
               chips: Sequence[str] = (), note: str = "", subtitle: str = ""):
        """One screenshot, full content width, with optional caption/chips/footnote."""
        sl, st = self._new()
        self._header(sl, st, title, subtitle)
        y = self._top()
        if caption:
            y = self._para(sl, st, self.ML, y, self.CW, caption, size=SZ["body"],
                           color=TEXT_BODY, spacing=1.22, max_h=1.0,
                           label="caption") + 0.20
        if chips:
            y = self._chips(sl, st, self.ML, y, chips) + 0.24
        bottom = 6.72 if note else BODY_BOTTOM
        self._picture(sl, st, image, self.ML, y, self.CW, bottom - y)
        if note:
            self._para(sl, st, self.ML, FOOTNOTE_Y, self.CW, note,
                       size=SZ["footnote"], color=TEXT_MUTED, spacing=1.12,
                       max_h=0.5, label="note")

    def split(self, title: str, image: str, *, card_title: str, para: str = "",
              bullets: Sequence[str] = (), chips: Sequence[str] = (),
              subtitle: str = ""):
        """Screenshot left, explanatory card right. The workhorse layout."""
        sl, st = self._new()
        self._header(sl, st, title, subtitle)
        top = self._top()
        bottom = 6.90
        img_w = self.CW * 0.583
        self._picture(sl, st, image, self.ML, top, img_w, bottom - top)
        cx = self.ML + img_w + 0.38
        cw = self.CW - img_w - 0.38
        self._rect(sl, st, cx, top, cw, bottom - top,
                   fill=SURFACE if self.style == "report" else SURFACE_ALT,
                   label="split-card")
        inner = cw - 0.84
        y = self._para(sl, st, cx + 0.42, top + 0.43, inner, card_title,
                       size=SZ["card_title"], bold=True, color=ACCENT,
                       spacing=1.08, max_h=1.0, label="split-title")
        if para:
            y = self._para(sl, st, cx + 0.42, y + 0.30, inner, para,
                           size=SZ["body"] - 0.5, color=TEXT_BODY, spacing=1.22,
                           max_h=1.9, label="split-para")
        if chips:
            y = self._chips(sl, st, cx + 0.42, y + 0.26, chips,
                            size=SZ["chip"] - 0.5, width=inner)
        if bullets:
            self._bullets(sl, st, cx + 0.42, y + 0.30, inner, bullets,
                          size=SZ["caption"], gap=7)

    def two_up(self, title: str, left: str, right: str, *, caption: str = "",
               left_caption: str = "", right_caption: str = "", subtitle: str = ""):
        """Two images side by side with captions underneath."""
        sl, st = self._new()
        self._header(sl, st, title, subtitle)
        y = self._top()
        if caption:
            y = self._para(sl, st, self.ML, y, self.CW, caption, size=SZ["body"],
                           color=TEXT_BODY, spacing=1.22, max_h=0.9,
                           label="caption") + 0.22
        cw = (self.CW - GUTTER) / 2
        cap_h = 0.42 if (left_caption or right_caption) else 0.0
        ih = 6.98 - cap_h - y
        self._picture(sl, st, left, self.ML, y, cw, ih, label="left")
        self._picture(sl, st, right, self.ML + cw + GUTTER, y, cw, ih, label="right")
        if cap_h:
            cy = y + ih + 0.10
            for cx, cap in ((self.ML, left_caption), (self.ML + cw + GUTTER, right_caption)):
                if cap:
                    self._para(sl, st, cx, cy, cw, cap, size=SZ["caption"], bold=True,
                               color=TEXT_PRIMARY, spacing=1.0, align=PP_ALIGN.CENTER,
                               max_h=0.4, label="imgcap")

    def flow(self, title: str, nodes: Sequence[tuple[str, str]], *, caption: str = "",
             fan_out: Sequence[tuple[str, str]] = (),
             cards: Sequence[tuple[str, Sequence[str]]] = (), subtitle: str = ""):
        """Left-to-right node diagram, optional fan-out, optional cards below."""
        sl, st = self._new()
        self._header(sl, st, title, subtitle)
        y = self._top()
        if caption:
            y = self._para(sl, st, self.ML, y, self.CW, caption, size=SZ["body"],
                           color=TEXT_BODY, spacing=1.22, max_h=0.8,
                           label="caption") + 0.22
        card_top = 4.72 if cards else BODY_BOTTOM
        region_h = card_top - 0.28 - y
        nh, ngap = 0.73, 0.13
        k = max(len(fan_out), 1)
        stack_h = k * nh + (k - 1) * ngap
        stack_top = y + max((region_h - stack_h) / 2, 0)
        mid = stack_top + stack_h / 2

        chain_w = 2.85
        surface = SURFACE if self.style == "report" else SURFACE_ALT
        for i, (label, sub) in enumerate(nodes):
            x = self.ML + i * (chain_w + 0.55)
            self._node(sl, st, x, mid - 0.525, chain_w, 1.05, label, sub,
                       fill=ACCENT_SOFT if i == len(nodes) - 1 and fan_out else surface)
            if i:
                self._connector(sl, x - 0.55, mid, x, mid)
        last_right = self.ML + (len(nodes) - 1) * (chain_w + 0.55) + chain_w
        fx = max(last_right + 0.55, SLIDE_W - MARGIN - 2.95)
        for j, (label, sub) in enumerate(fan_out):
            ny = stack_top + j * (nh + ngap)
            self._node(sl, st, fx, ny, 2.95, nh, label, sub, fill=surface)
            self._connector(sl, last_right, mid, fx, ny + nh / 2)
        if cards:
            ch = BODY_BOTTOM - card_top + 0.18
            cw = (self.CW - GUTTER) / len(cards)
            for i, (ct, items) in enumerate(cards):
                cx = self.ML + i * (cw + GUTTER)
                self._rect(sl, st, cx, card_top, cw, ch, fill=surface,
                           label=f"flowcard{i}")
                self._para(sl, st, cx + 0.40, card_top + 0.28, cw - 0.8, ct,
                           size=SZ["card_title"] - 2, bold=True, color=ACCENT,
                           spacing=1.0, max_h=0.5, label="flowcard-title")
                self._bullets(sl, st, cx + 0.40, card_top + 0.85, cw - 0.8, items,
                              size=SZ["caption"], gap=7)

    def _node(self, sl, st, x, y, w, h, label, sub, fill=SURFACE):
        self._rect(sl, st, x, y, w, h, fill=fill, label=f"node:{label}")
        self._text(sl, st, x + 0.14, y + h * 0.17, w - 0.28, 0.34,
                   [{"text": label, "size": SZ["caption"], "bold": True,
                     "color": TEXT_PRIMARY, "spacing": 1.0}],
                   align=PP_ALIGN.CENTER, check=False)
        if sub:
            self._text(sl, st, x + 0.14, y + h * 0.55, w - 0.28, 0.30,
                       [{"text": sub, "size": SZ["footnote"] - 1, "bold": False,
                         "color": TEXT_MUTED, "spacing": 1.0}],
                       align=PP_ALIGN.CENTER, check=False)

    def _connector(self, sl, x1, y1, x2, y2):
        c = sl.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        c.line.color.rgb = ACCENT_LIGHT
        c.line.width = Pt(1.5)
        self._colors_used.add(str(ACCENT_LIGHT))

    def rows(self, title: str, entries: Sequence[tuple[str, str]], *, intro: str = "",
             subtitle: str = ""):
        """Label/detail rows — findings, observations, spec tables. Max 5."""
        sl, st = self._new()
        self._header(sl, st, title, subtitle)
        y = self._top()
        if intro:
            y = self._para(sl, st, self.ML, y, self.CW, intro, size=SZ["body"],
                           color=TEXT_BODY, spacing=1.2, max_h=0.7,
                           label="intro") + 0.24
        n = len(entries)
        if n > 5:
            self._note(st, "row-count", f"{n} rows given; 5 is the maximum that fits")
        gap = 0.12
        rh = (BODY_BOTTOM - y - gap * (n - 1)) / max(n, 1)
        lw = min(3.5, self.CW * 0.30)
        for i, (label, detail) in enumerate(entries):
            ry = y + i * (rh + gap)
            self._rect(sl, st, self.ML, ry, self.CW, rh,
                       fill=SURFACE if self.style == "report" else SURFACE_ALT,
                       label=f"row{i}")
            self._rect(sl, st, self.ML, ry, 0.075, rh, fill=ACCENT, register=False)
            self._para(sl, st, self.ML + 0.33, ry + 0.20, lw, label,
                       size=SZ["caption"] + 1.5, bold=True, color=TEXT_PRIMARY,
                       spacing=1.05, max_h=rh - 0.3, label=f"rowl{i}")
            dx = self.ML + 0.33 + lw + 0.30
            self._para(sl, st, dx, ry + 0.20, SLIDE_W - MARGIN - dx, detail,
                       size=SZ["caption"], color=TEXT_BODY, spacing=1.18,
                       max_h=rh - 0.3, label=f"rowd{i}")

    def closing(self, *, title: str = "Thank You", para: str = "",
                summary: Sequence[str] = (), contact: Sequence[str] = (),
                footer: str = "www.overview.ai"):
        """Closing slide: navy field, logo, optional summary and contact panels."""
        sl, st = self._new(dark=True)
        self._logo(sl, st, 0.95, 0.72, 2.35, on_dark=True)
        y = self._para(sl, st, 0.95, 1.85, 9.0, title, size=SZ["display"] + 6, bold=True,
                       color=TEXT_ON_DARK, spacing=1.0, max_h=1.2, label="thanks")
        if para:
            y = self._para(sl, st, 0.95, y + 0.28, 9.6, para, size=SZ["lead"],
                           color=ACCENT_ON_DARK, spacing=1.25, max_h=1.1,
                           label="closing-para")
        panels = [p for p in ((summary, "Summary"), (contact, "Contact")) if p[0]]
        py = max(y + 0.42, 4.10)
        ph = min(6.55 - py, 2.6)
        pw = 5.6
        for i, (lines, heading) in enumerate(panels):
            px = 0.95 + i * (pw + 0.5)
            self._rect(sl, st, px, py, pw, ph, fill=DARK_BG_ALT, label=f"panel{i}")
            self._para(sl, st, px + 0.45, py + 0.30, pw - 0.9, heading,
                       size=SZ["card_title"] - 1, bold=True, color=TEXT_ON_DARK,
                       spacing=1.0, max_h=0.4, label=f"panelh{i}")
            for j, line in enumerate(lines[:4]):
                self._para(sl, st, px + 0.45, py + 0.88 + j * 0.44, pw - 0.9, line,
                           size=SZ["body"] - 0.5, color=TEXT_ON_DARK, spacing=1.0,
                           max_h=0.4, label=f"panel{i}l{j}")
        if footer:
            self._para(sl, st, 0.95, 6.86, 4.0, footer, size=SZ["footnote"] - 0.5,
                       color=TEXT_ON_DARK, spacing=1.0, label="footer")

    # ---------------- validation ----------------

    def check(self) -> list[Issue]:
        """Run every structural and brand rule. Called automatically by save()."""
        issues = list(self.issues)
        if self._font_note:
            issues.append(Issue(0, "font-fallback", self._font_note, "warning"))

        for st in self.states:
            texts = [r for r in st.rects if r.kind == "text"]
            images = [r for r in st.rects if r.kind in ("image", "logo")]
            for i, a in enumerate(texts):
                for b in texts[i + 1:]:
                    if a.overlap_area(b) > 0.02:
                        issues.append(Issue(st.index, "collision",
                                            f"text '{a.label}' overlaps text '{b.label}' "
                                            f"({a.overlap_area(b):.2f} sq in)"))
            for t in texts:
                for im in images:
                    if t.overlap_area(im) > 0.02:
                        issues.append(Issue(st.index, "collision",
                                            f"text '{t.label}' overlaps image '{im.label}'"))
            for i, a in enumerate(images):
                for b in images[i + 1:]:
                    if a.overlap_area(b) > 0.05 and "disc" not in (a.label + b.label):
                        issues.append(Issue(st.index, "collision",
                                            f"image '{a.label}' overlaps image '{b.label}'"))

        allowed = {a.upper() for a in ALLOWED_COLORS}
        for c in sorted(self._colors_used):
            if c.upper() not in allowed:
                issues.append(Issue(0, "off-brand-colour",
                                    f"#{c} is not in the Overview palette"))

        if self.states:
            if not self.states[0].has_logo:
                issues.append(Issue(1, "missing-logo",
                                    "the opening slide must carry the logo"))
            # When the deck ends with transplanted boilerplate, the real
            # closing slide is the template's own and carries the template's
            # logo — this engine never laid it out and cannot see it. Checking
            # the last AUTHORED slide instead would flag whatever happens to
            # precede the boilerplate, which is not the closing slide at all.
            if not getattr(self, "_ends_with_template", False) and not self.states[-1].has_logo:
                issues.append(Issue(len(self.states), "missing-logo",
                                    "the closing slide must carry the logo"))
        return issues

    def skeleton_slide(
        self,
        name: str,
        *,
        image: str | Path | None = None,
        tokens: dict | None = None,
    ) -> None:
        """Place a boilerplate slide from the skill's owned skeletons —
        library, capabilities, team, thank-you — exactly as designed.

        Each skeleton is a single-slide pptx extracted once from the company
        template (assets/skeletons/, see template_slides.py). `image` fills
        the slide's "Insert screenshot here" hole (the library slide has
        one); `tokens` fills {{token}} text holes if a skeleton ever grows
        any. Unfilled holes raise rather than ship placeholder text.

        Slides added this way bypass check() by design — they are known-good
        company content, not something this engine laid out, and measuring
        them against its capacities would report failures it cannot fix.
        If a filled skeleton looks wrong in the render, fix the skeleton
        (a maintainer change, reapplied on re-extraction) or shorten the
        content — never patch the built deck.
        """
        from template_slides import append

        append(self.prs, name, image=image, tokens=tokens)
        self._template_slides = getattr(self, "_template_slides", 0) + 1
        # Position in the finished file, so the audit can tell what this
        # engine laid out from what it merely carried over.
        self._template_indices = getattr(self, "_template_indices", [])
        self._template_indices.append(len(self.prs.slides._sldIdLst))
        self._ends_with_template = True

    def template_slide(self, name: str, *, image: str | Path | None = None) -> None:
        """Deprecated alias for skeleton_slide()."""
        self.skeleton_slide(name, image=image)

    def save(self, path: str | Path | None = None) -> Path:
        out = Path(path or self.out_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        issues = self.check()
        errors = [i for i in issues if i.severity == "error"]
        for i in issues:
            print(i)
        if errors and self.strict:
            raise LayoutError(
                f"{len(errors)} layout/brand error(s); deck not written. "
                "Fix the content or pick a layout with more room — never nudge "
                "coordinates by hand."
            )
        idx = getattr(self, "_template_indices", [])
        if idx:
            # Self-describing: the audit reads this instead of guessing which
            # slides came from the reference template.
            self.prs.core_properties.keywords = (
                "ovdeck:template-slides=" + ",".join(str(i) for i in idx)
            )
        self.prs.save(str(out))
        # the circular title image is embedded at save time; the scratch copy
        # on disk is no longer referenced
        tmp = out.parent / "_title_disc.png"
        if tmp.exists():
            try:
                tmp.unlink()
            except OSError:
                pass
        n_tmpl = getattr(self, "_template_slides", 0)
        n_all = len(self.states) + n_tmpl
        carried = f" (+{n_tmpl} carried from template)" if n_tmpl else ""
        print(f"\nsaved: {out}  ({n_all} slides{carried}, {len(errors)} error(s), "
              f"{len(issues) - len(errors)} warning(s), style={self.style}, "
              f"font={self.font})")
        return out


def _circle_crop(src: str, dst: Path, size: int = 1000) -> Path:
    """Centre-crop to a square and mask to a circle (title-slide motif)."""
    from PIL import ImageDraw

    dst.parent.mkdir(parents=True, exist_ok=True)
    im = Image.open(src).convert("RGB")
    w, h = im.size
    s = min(w, h)
    im = im.crop(((w - s) // 2, (h - s) // 2, (w + s) // 2, (h + s) // 2))
    im = im.resize((size, size), Image.LANCZOS)
    mask = Image.new("L", (size, size), 0)
    ImageDraw.Draw(mask).ellipse((0, 0, size, size), fill=255)
    im.putalpha(mask)
    im.save(dst)
    return dst


__all__ = ["Deck", "LayoutError", "Issue", "text_width_in", "text_height_in",
           "wrap_lines", "font_installed", "TOKENS"]
