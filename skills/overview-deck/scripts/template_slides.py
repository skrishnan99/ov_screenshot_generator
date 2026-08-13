#!/usr/bin/env python3
"""Boilerplate slides as owned skeletons: extract once, fill, place exactly.

The library, capability and team slides are standing company content. They
used to be re-authored through the layout engine (a different approximation
every deck) and then transplanted verbatim from the 26 MB reference template
(byte-exact, but unownable — its defects shipped into every deck and could
never be fixed without editing the shared example).

This module implements the third way, borrowed from the plugin's skeleton
pipeline: each boilerplate slide is EXTRACTED ONCE into a small single-slide
pptx under assets/skeletons/, which the skill then owns. A sidecar YAML with
the same stem describes its content holes. Build scripts query a skeleton,
fill its holes (the library screenshot; {{tokens}} if any are ever added) and
append it verbatim — and because the skeletons are ours, a layout defect is
fixed once in the skeleton file and stays fixed.

Maintainer operations:

    python template_slides.py                # list skeletons and their holes
    python template_slides.py --extract      # (re)build skeletons from the
                                             # reference template; run only
                                             # when the company template
                                             # changes, then re-apply fixes
                                             # and review the diff

Build-script use (via ovdeck):

    d.skeleton_slide("library", image=run/"deliverables/screenshots/12_library.png")
    for name in ("capabilities", "team", "thank_you"):
        d.skeleton_slide(name)
"""

from __future__ import annotations

import sys
from pathlib import Path

SKILL = Path(__file__).resolve().parent.parent
PLUGIN_ROOT = SKILL.parent.parent
TEMPLATE = SKILL / "assets" / "example-decks" / "Overview AI blank test report.pptx"
SKELETON_DIR = SKILL / "assets" / "skeletons"

# The plugin owns the OPC surgery and the skeleton introspection; this skill
# ships inside it. Same three-levels-up convention publish.py uses.
if str(PLUGIN_ROOT) not in sys.path:
    sys.path.insert(0, str(PLUGIN_ROOT))

# Where each boilerplate slide lives in the reference template (1-based, as
# the slides read in PowerPoint). Used by --extract only; at build time the
# skeleton files are the source of truth.
TEMPLATE_SLIDES = {
    "capabilities": 11,      # "5 factors make us unique"
    "defect_generator": 12,  # "Can't wait for the rare defect? Generate it."
    "integration": 13,       # "Integration built for everyone."
    "team": 14,              # Team & Locations
    "thank_you": 15,         # closing
}

# Content-hole skeletons imported from the v1 corpus (assets/reference/,
# copied verbatim from deck/skeletons and the case_study repo). Unlike the
# template boilerplate above, these carry {{token}} text holes and PICTURE /
# placeholder image slots, and are filled with the v1 semantics
# (fill_tokens + fill_images) that shipped hundreds of slides. They are
# normalised to the deck canvas at --extract exactly like the template
# slides — THE CANVAS INVARIANT applies to every imported skeleton.
CONTENT_SKELETONS = {
    "recipe_title_ov80i": "recipe_title_ov80i.pptx",
    "recipe_title_ov20i": "recipe_title_ov20i.pptx",
    "results_image": "results_image.pptx",
    "concise_results_classifier": "concise_results_classifier.pptx",
    "concise_results_segmenter": "concise_results_segmenter.pptx",
    # v1's library slide (one screenshot placeholder, no tokens). Replaces
    # the company template's slide 10 as the deck's library section; the
    # _fix_library_subtitle fixup still applies — same subtitle box.
    # ov20i/ov10i keep this one; ov80i gets its own variant below.
    "library": "library.pptx",
    # the OV80i library slide (updated copy, same one-placeholder shape);
    # selected by the spec's `when: {camera.variant: ov80i}`.
    "library_ov80i": "library_ov80i.pptx",
}
REFERENCE_DIR = SKILL / "assets" / "reference"

# The closing run every report carries, in the TEMPLATE's own order
# (slides 11-15). Library is not listed here: it is its own numbered section
# of the default deck, placed earlier with this run's screenshot.
DEFAULT_CLOSING = (
    "capabilities",
    "defect_generator",
    "integration",
    "team",
    "thank_you",
)


class TemplateError(RuntimeError):
    """A skeleton is missing or does not match what the build needs."""


def all_skeleton_names() -> set[str]:
    return set(TEMPLATE_SLIDES) | set(CONTENT_SKELETONS)


def skeleton_path(name: str) -> Path:
    if name not in all_skeleton_names():
        raise TemplateError(
            f"unknown skeleton {name!r}; have {sorted(all_skeleton_names())}"
        )
    p = SKELETON_DIR / f"{name}.pptx"
    if not p.exists():
        raise TemplateError(
            f"skeleton not found at {p}. Skeletons ship with the skill; "
            f"regenerate with `python {Path(__file__).name} --extract` if the "
            f"directory was lost."
        )
    return p


def profile(name: str) -> dict:
    """The skeleton's content holes, via the plugin's introspection: tokens,
    image slots with geometry, sidecar guidance, drift warnings."""
    from deck.slots import skeleton_profile

    return skeleton_profile(str(skeleton_path(name)))


def append(prs, name: str, image: str | Path | None = None,
           tokens: dict | None = None,
           images: list | None = None) -> None:
    """Fill skeleton `name` and append it to an open Presentation, verbatim.

    `images` (or the single-image `image` alias) fill the slide's image
    holes; `tokens` fill its {{token}} text holes. Unknown surplus is an
    error — a hole that silently stays unfilled ships placeholder text to a
    customer.

    Two fill modes, decided by the registry:
    - TEMPLATE_SLIDES: "Insert screenshot here" placeholders only; standing
      PICTUREs are content and never touched.
    - CONTENT_SKELETONS (imported v1 corpus): the v1 semantics — PICTURE
      slots are replaceable holes (largest first), then placeholders; fewer
      images than picture slots leaves the remainder standing (the recipe
      title keeps its hero photo when no image is passed).
    """
    from pptx import Presentation

    from deck.assemble import (
        append_slide,
        bake_theme_colors,
        fill_images,
        fill_tokens,
        find_tokens,
        image_slots,
    )

    if images is None:
        images = [image] if image is not None else []
    if name == "thank_you":
        # The contact block is tokenized ({{contact_*}}) and defaults to the
        # engineer profile — merged HERE, the lowest fill point, so every
        # caller (the closing run, hand-written v1 build scripts) gets the
        # right contact without passing anything. Missing profile fields
        # resolve to visibly generic placeholders, never to a wrong person.
        from core.engineer import load_profile

        contact, _ = load_profile()
        tokens = {
            "contact_name": contact["name"],
            "contact_email": contact["email"],
            "contact_phone": contact["phone"],
            **(tokens or {}),
        }
    src = Presentation(str(skeleton_path(name)))
    slide = src.slides[0]

    # THE CANVAS INVARIANT (see the extraction section): a skeleton on a
    # different canvas than the deck renders shrunken in a corner. Refuse
    # loudly — the fix is re-extraction, not a build-time workaround.
    if (src.slide_width, src.slide_height) != (prs.slide_width, prs.slide_height):
        raise TemplateError(
            f"skeleton {name!r} canvas is "
            f"{src.slide_width}x{src.slide_height} EMU but the deck is "
            f"{prs.slide_width}x{prs.slide_height}. The skeletons predate a "
            f"canvas change — regenerate with `python {Path(__file__).name} "
            f"--extract`."
        )

    if tokens:
        have = set(find_tokens(slide))
        surplus = set(tokens) - have
        if surplus:
            raise TemplateError(
                f"skeleton {name!r} has no {{{{token}}}} hole named "
                f"{sorted(surplus)}; it has {sorted(have) or 'none'}"
            )
        fill_tokens(slide, tokens)
    leftover = find_tokens(slide)
    if leftover:
        raise TemplateError(
            f"skeleton {name!r} still has unfilled tokens {leftover}; "
            f"pass them via tokens={{...}}"
        )

    for img in images:
        if not Path(img).exists():
            raise TemplateError(f"image for skeleton {name!r} not found: {img}")

    if name in CONTENT_SKELETONS:
        # v1 semantics: pictures are holes (largest first), then any
        # placeholders. Fewer images than picture slots leaves the remainder
        # standing; leftover PLACEHOLDERS would ship their note, so refuse.
        from deck.assemble import MSO_PICTURE

        slots = image_slots(slide)
        n_placeholders = sum(1 for sh in slots if sh.shape_type != MSO_PICTURE)
        n_pics = len(slots) - n_placeholders
        if len(images) > len(slots):
            raise TemplateError(
                f"skeleton {name!r} has {len(slots)} image slot(s); got {len(images)}"
            )
        if n_placeholders and len(images) < len(slots):
            raise TemplateError(
                f"skeleton {name!r} has {n_placeholders} placeholder hole(s) "
                f"after {n_pics} picture slot(s); {len(images)} image(s) would "
                f"leave placeholder text on the slide"
            )
        fill_images(slide, [str(i) for i in images])
    else:
        # Template boilerplate: placeholders only; standing PICTUREs are
        # content and must never be replaced. (fill_images targets pictures
        # first, so it is the wrong tool here.)
        from deck.assemble import MSO_PICTURE, _fill_placeholder

        holes = [sh for sh in image_slots(slide) if sh.shape_type != MSO_PICTURE]
        for img in images:
            if not holes:
                raise TemplateError(
                    f"skeleton {name!r} has no 'Insert screenshot here' hole "
                    f"to put an image in"
                )
            _fill_placeholder(slide, holes.pop(0), Path(img))
        if holes:
            # An unfilled hole would ship its placeholder note to a customer.
            raise TemplateError(
                f"skeleton {name!r} has {len(holes)} unfilled screenshot hole(s); "
                f"pass images=..."
            )

    # Idempotent when the extraction already baked; load-bearing if a skeleton
    # was hand-edited afterwards and picked up scheme colours again.
    bake_theme_colors(slide)
    ctx = {"used": {str(p.partname) for p in prs.part.package.iter_parts()},
           "by_hash": {}}
    append_slide(prs, src, ctx)


# ---------------------------------------------------------------- extraction
#
# Owning the skeletons is the point: template defects are fixed HERE, once,
# and re-extraction reapplies them. Never fix a defect in the built deck.
#
# THE CANVAS INVARIANT. The deprecated deck generator was reliable because
# every slide shared one coordinate space: its base presentation WAS the
# first skeleton, so the canvas came from the skeletons themselves. ovdeck
# authors at 13.333x7.5in while the company template is 10x5.625in — same
# 16:9, different absolute EMUs — so a verbatim transplant huddled in the
# top-left 75% of the deck. Extraction therefore normalises every skeleton
# to ovdeck's canvas, scaling geometry AND type by the exact ratio (4/3),
# once, here, where a maintainer can eyeball the result. Build-time append
# stays verbatim and refuses a skeleton whose canvas disagrees with the deck.


def _scale_part_xml(blob: bytes, factor: float) -> bytes:
    """Uniformly scale one slide/layout/master part's XML.

    Both canvases are 16:9, so a single factor maps everything: shape frames
    (off/ext, and group child spaces chOff/chExt — scaling all four uniformly
    is self-consistent), font sizes (sz, 1/100 pt), line widths, text insets,
    fixed paragraph spacing, table geometry and shadow offsets. Percentages
    (srcRect, spcPct, normAutofit) need no scaling and are left alone; the
    isdigit guard also skips non-numeric attrs like p:ph's sz="quarter".
    """
    from lxml import etree

    root = etree.fromstring(blob)

    def scale(el, *names):
        for name in names:
            v = el.get(name)
            if v is not None and v.lstrip("-").isdigit():
                el.set(name, str(round(int(v) * factor)))

    for el in root.iter():
        tag = etree.QName(el).localname
        if tag in ("off", "chOff"):
            scale(el, "x", "y")
        elif tag in ("ext", "chExt"):
            scale(el, "cx", "cy")
        elif tag == "ln":
            scale(el, "w")
        elif tag in ("rPr", "defRPr", "endParaRPr"):
            scale(el, "sz")
        elif tag == "bodyPr":
            scale(el, "lIns", "tIns", "rIns", "bIns")
        elif tag in ("spcPts", "buSzPts"):
            scale(el, "val")
        elif tag == "gridCol":
            scale(el, "w")
        elif tag == "tr":
            scale(el, "h")
        elif tag in ("outerShdw", "innerShdw", "prstShdw"):
            scale(el, "blurRad", "dist")
    return etree.tostring(
        root, xml_declaration=True, encoding="UTF-8", standalone=True
    )

def _strip_page_numbers(slide) -> None:
    """The template's slides carry hardcoded "NN / 15" page markers. Any
    number would be wrong in a built deck, whose slide count varies."""
    import re

    from deck.assemble import iter_shapes

    page = re.compile(r"^\s*\d{1,2}\s*/\s*\d{1,2}\s*$")
    for sh in list(iter_shapes(slide)):
        if sh.has_text_frame and page.match(sh.text_frame.text or ""):
            sh._element.getparent().remove(sh._element)


def _fix_library_subtitle(slide) -> None:
    """The subtitle box is one line high (H 0.24) for a line that is marginal
    at its width: with substituted fonts it sometimes wraps, and the overflow
    lands on the screenshot at T 1.66. Give the box two lines of room — a
    wrapped second line then ends at 1.59, clear of the image."""
    from pptx.util import Inches

    from deck.assemble import iter_shapes

    for sh in iter_shapes(slide):
        if sh.has_text_frame and sh.text_frame.text.startswith("Easier root cause"):
            sh.height = Inches(0.45)


def _tokenize_thank_you_contact(slide) -> None:
    """Rewrite the contact block's literal runs into {{contact_*}} tokens so
    the slide carries the SE who actually ran the visit, not whoever was on
    the template. Run at extraction, where a template change that breaks the
    pattern should fail loudly at the maintainer, never at build time.

    Targeting is by pattern, not by name: the email run contains "@", the
    phone run is digits/punctuation, and the one remaining non-"Thank you"
    run is the name. Only run text changes — formatting stays."""
    import re as _re

    from deck.assemble import iter_shapes

    phone_re = _re.compile(r"^[\d\s()+\-.]{7,}$")
    email_run = phone_run = name_run = None
    leftovers = []
    for sh in iter_shapes(slide):
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                text = (run.text or "").strip()
                if not text or "thank" in text.lower():
                    continue
                if "@" in text and " " not in text:
                    email_run = run
                elif phone_re.match(text):
                    phone_run = run
                else:
                    leftovers.append(run)
    if len(leftovers) == 1:
        name_run = leftovers[0]
    if not (email_run and phone_run and name_run):
        raise TemplateError(
            "thank_you: could not identify the contact runs to tokenize "
            f"(email={bool(email_run)}, phone={bool(phone_run)}, "
            f"name candidates={len(leftovers)}). The template's contact "
            "block changed — update _tokenize_thank_you_contact."
        )
    name_run.text = "{{contact_name}}"
    email_run.text = "{{contact_email}}"
    phone_run.text = "{{contact_phone}}"


def _strip_integration_page_number(slide) -> None:
    """The integration slide carries its page number as a bare "12", which
    the NN / NN pattern cannot catch — and a generic bare-number strip would
    also delete legitimate design numerals (defect_generator's 01/02/03 step
    markers, its 12.4x stat). Target exactly this one box."""
    from deck.assemble import iter_shapes

    for sh in list(iter_shapes(slide)):
        if sh.has_text_frame and sh.text_frame.text.strip() == "12":
            sh._element.getparent().remove(sh._element)


# Applied in order after transplant, before save. _strip_page_numbers runs on
# every skeleton; per-name fixups follow.
FIXUPS: dict[str, tuple] = {
    "*": (_strip_page_numbers,),
    "library": (_fix_library_subtitle,),
    # same "Easier root cause..." subtitle box, same wrap hazard
    "library_ov80i": (_fix_library_subtitle,),
    "integration": (_strip_integration_page_number,),
    "thank_you": (_tokenize_thank_you_contact,),
}


def _shrink_pictures(slide, scale: float = 1.0) -> int:
    """Re-embed each picture at the size its frame actually displays.

    The template's slides carry photography at capture resolution — one slide
    weighed 16 MB for four pictures displayed a few inches wide. The raster is
    downscaled to the frame size at the plugin's EMBED_DPI and re-encoded
    (PNG/JPEG chosen empirically, transparency preserved); only the blip's
    image part is re-pointed, so geometry, crops and stretch stay byte-exact.
    Uniform scaling keeps any srcRect crop fractions valid. Failures leave the
    original picture — a bigger skeleton beats a broken one.
    """
    import io

    from PIL import Image
    from pptx.oxml.ns import qn
    from pptx.util import Emu

    from deck.assemble import EMBED_DPI, MSO_PICTURE, _encode_smaller, iter_shapes

    shrunk = 0
    for sh in iter_shapes(slide):
        if sh.shape_type != MSO_PICTURE:
            continue
        try:
            blip = sh._element.blipFill.find(qn("a:blip"))
            part = sh.part.related_part(blip.get(qn("r:embed")))
            with Image.open(io.BytesIO(part.blob)) as im:
                tw = max(1, int(Emu(sh.width).inches * EMBED_DPI * scale))
                th = max(1, int(Emu(sh.height).inches * EMBED_DPI * scale))
                if im.width <= tw and im.height <= th:
                    continue
                # NB: never name this `scale` — it would shadow the
                # parameter and compound across pictures (each raster ends up
                # shrunk by the PREVIOUS picture's resize ratio, collapsing
                # geometrically down the slide: 687 -> 180 -> 47 -> 12 px).
                ratio = min(tw / im.width, th / im.height)
                if im.mode == "P":
                    out = im.convert("RGBA" if "transparency" in im.info else "RGB")
                elif im.mode == "CMYK":
                    out = im.convert("RGB")
                else:
                    out = im.copy()
                out = out.resize(
                    (max(1, int(im.width * ratio)), max(1, int(im.height * ratio))),
                    Image.LANCZOS,
                )
                tmp = _encode_smaller(out)
            _, new_rId = sh.part.get_or_add_image_part(tmp)
            blip.set(qn("r:embed"), new_rId)
            shrunk += 1
        except Exception:
            continue
    if shrunk:
        _drop_orphan_image_rels(slide)
    return shrunk


def _drop_orphan_image_rels(slide) -> None:
    """Re-pointing a blip leaves the old image relationship on the slide
    part, and the transplant copies every relationship — so without this the
    skeleton carries the original AND the shrunk raster and comes out BIGGER
    than before. Drop image rels nothing in the slide XML references any
    more. The same part may back several pictures, so every referenced rId is
    collected before anything is dropped."""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.oxml.ns import qn

    used = set()
    for el in slide._element.iter():
        for attr in (qn("r:embed"), qn("r:link")):
            v = el.get(attr)
            if v:
                used.add(v)
    part = slide.part
    for rId, rel in list(part.rels.items()):
        if rel.reltype == RT.IMAGE and rId not in used:
            part.drop_rel(rId)


def extract(template: Path | None = None, out_dir: Path | None = None) -> list[Path]:
    """Lift each boilerplate slide out of the reference template into a small
    single-slide skeleton. Maintainer operation: rerun when the company
    template changes, then re-apply any skeleton fixes and review the diff.

    Each skeleton carries only the parts its slide references (its layout,
    master, theme and media travel via the plugin's OPC import), so the files
    are small and fully self-contained. Theme colours are baked at extraction,
    making the skeleton immune to recolouring by any host deck.
    """
    from pptx import Presentation

    from deck.assemble import append_slide, bake_theme_colors

    src_path = Path(template) if template else TEMPLATE
    if not src_path.exists():
        raise TemplateError(f"reference template not found at {src_path}")
    out = Path(out_dir) if out_dir else SKELETON_DIR
    out.mkdir(parents=True, exist_ok=True)

    src = Presentation(str(src_path))

    # Normalise to the deck's canvas (see THE CANVAS INVARIANT above). Both
    # must be the same aspect or a single factor cannot map the geometry.
    from pptx.util import Inches

    from ovdeck import SLIDE_H, SLIDE_W

    factor = Inches(SLIDE_W) / src.slide_width
    if abs(factor - Inches(SLIDE_H) / src.slide_height) > 1e-4:
        raise TemplateError(
            f"template canvas {src.slide_width}x{src.slide_height} EMU is a "
            f"different aspect than the deck's {SLIDE_W}x{SLIDE_H} in; a "
            f"uniform scale cannot map it"
        )

    written = []
    for name, num in TEMPLATE_SLIDES.items():
        if num - 1 >= len(src.slides):
            raise TemplateError(
                f"template has {len(src.slides)} slides; {name!r} expected at "
                f"{num}. The template changed — update TEMPLATE_SLIDES."
            )
        written.append(_import_one(src, num - 1, name, factor, out))

    # Content skeletons: single-slide files from the v1 corpus, same
    # normalisation pipeline (shrink, fixups, bake, transplant, scale).
    # Their sidecar YAMLs travel with them so profile() keeps working.
    for name, fname in CONTENT_SKELETONS.items():
        ref = REFERENCE_DIR / fname
        if not ref.exists():
            raise TemplateError(f"reference skeleton missing: {ref}")
        ref_pres = Presentation(str(ref))
        f2 = Inches(SLIDE_W) / ref_pres.slide_width
        if abs(f2 - Inches(SLIDE_H) / ref_pres.slide_height) > 1e-4:
            raise TemplateError(f"{fname}: aspect differs from the deck canvas")
        written.append(_import_one(ref_pres, 0, name, f2, out))
        sidecar = ref.with_suffix(".yaml")
        if sidecar.exists():
            (out / f"{name}.yaml").write_text(sidecar.read_text())
    return written


def _import_one(src_pres, index: int, name: str, factor: float, out: Path) -> Path:
    """Normalise one source slide into skeletons/<name>.pptx on the deck
    canvas: shrink rasters to post-scale display size, apply fixups, bake
    theme colours, transplant, scale every imported part (slide, layout AND
    master — layouts/masters carry positioned chrome of their own)."""
    from pptx import Presentation
    from pptx.util import Inches

    from deck.assemble import append_slide, bake_theme_colors

    from ovdeck import SLIDE_H, SLIDE_W

    # A fresh base per skeleton, on the DECK's canvas. The default
    # python-pptx master stays in the file unused (~10 KB) — harmless.
    base = Presentation()
    base.slide_width = Inches(SLIDE_W)
    base.slide_height = Inches(SLIDE_H)
    slide = src_pres.slides[index]
    # Shrink targets the size the picture will DISPLAY at post-scale, so
    # the raster still meets EMBED_DPI on the larger canvas.
    n = _shrink_pictures(slide, scale=factor)
    # Fixups run on the source slide: the transplanted copy's part is a
    # raw OPC Part without python-pptx's typed shape API. Mutating the
    # in-memory source is safe — it is never saved back.
    for fix in FIXUPS.get("*", ()) + FIXUPS.get(name, ()):
        fix(slide)
    bake_theme_colors(slide)
    before = {str(p.partname) for p in base.part.package.iter_parts()}
    ctx = {"used": set(before), "by_hash": {}}
    append_slide(base, src_pres, ctx, index=index)
    for part in base.part.package.iter_parts():
        if str(part.partname) in before:
            continue
        if part.content_type.endswith(
            ("slide+xml", "slideLayout+xml", "slideMaster+xml")
        ):
            part._blob = _scale_part_xml(part.blob, factor)
    dest = out / f"{name}.pptx"
    base.save(str(dest))
    print(f"  {dest.relative_to(SKILL)}  {dest.stat().st_size / 1024:.0f} KB"
          + (f"  ({n} picture(s) re-embedded at display size)" if n else ""))
    return dest


def main() -> int:
    if "--extract" in sys.argv:
        extract()
        return 0
    print(f"skeletons in {SKELETON_DIR.relative_to(SKILL)}:\n")
    for name in TEMPLATE_SLIDES:
        try:
            p = profile(name)
        except TemplateError as e:
            print(f"  {name:17} MISSING ({e})")
            continue
        holes = [s["name"] + ("*" if not s["is_picture"] else "")
                 for s in p["slots"]]
        star = " [default]" if name in DEFAULT_CLOSING or name == "library" else ""
        print(f"  {name:17} tokens={p['tokens'] or '—'}  "
              f"image-holes={holes or '—'}  {p['title'][:40]!r}{star}")
        for w in p["warnings"]:
            print(f"        drift: {w}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
